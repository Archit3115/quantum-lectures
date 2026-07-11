from pptx import Presentation
from pptx.util import Emu, Pt
import math, sys
_DEF={"session1":"session1_quantum_clustering_qiskit.pptx","session2":"session2_quantum_vision_qiskit.pptx"}
_arg=sys.argv[1] if len(sys.argv)>1 else "session1"
OUT="/Users/sentry/Work/Lectures/NIT Warangal/"+_DEF.get(_arg,_arg)
p=Presentation(OUT); SW,SH=p.slide_width,p.slide_height
n=len(p.slides._sldIdLst)
print("SLIDES:",n)
BADWORDS=("Lorem","Ipsum","Your name here","Click to add","supplied to Qiskit")
oob=[]; lorem=[]; imgs=0; overflow=[]
def est_height_in(shape):
    try: w_in=shape.width/914400
    except: return 0
    if w_in<=0: return 0
    total=0.0
    for para in shape.text_frame.paragraphs:
        txt="".join(r.text for r in para.runs)
        sz=None
        for r in para.runs:
            if r.font.size: sz=r.font.size.pt; break
        sz=sz or 12
        cw=(sz*0.5)/72.0
        cpl=max(1,int((w_in-0.1)/cw))
        lines=max(1,math.ceil(len(txt)/cpl)) if txt else 1
        sa=(para.space_after.pt if para.space_after else 3)
        total+= lines*(sz*1.18/72.0) + sa/72.0
    return total
for i,s in enumerate(p.slides,1):
    alltext=""
    for sh in s.shapes:
        if sh.shape_type==13: imgs+=1
        try: alltext+=" "+(sh.text_frame.text or "")
        except: pass
        try:
            l,t,w,h=sh.left,sh.top,sh.width,sh.height
            if None not in (l,t,w,h):
                if l<-9144 or t<-9144 or l+w>SW+9144 or t+h>SH+18288:
                    oob.append((i,round((l+w)/914400,1),round((t+h)/914400,1)))
        except: pass
        try:
            if sh.has_text_frame and sh.shape_type not in (13,19) and (sh.text_frame.text or "").strip():
                eh=est_height_in(sh); bh=sh.height/914400
                if eh>bh+0.25:
                    overflow.append((i,round(eh,1),round(bh,1),(sh.text_frame.text or "")[:30]))
        except: pass
    for bw in BADWORDS:
        if bw in alltext: lorem.append((i,bw))
print("images embedded:",imgs)
print("leftover template text:",lorem)
print("OOB (decorative bleeds expected):",len(oob))
print("TEXT OVERFLOW candidates:",len(overflow))
for o in overflow: print("   slide",o[0],"est",o[1],"vs box",o[2],"|",o[3])

# ---- strict THIRD-PERSON voice audit (content rule #3) ----
import re
FORBID=r"\b(you|your|you're|yours|we|we're|our|ours|us|my|me|i'm|let's|lets)\b"
voice=[]
for i,s in enumerate(p.slides,1):
    for sh in s.shapes:
        try: t=sh.text_frame.text or ""
        except: continue
        t=re.sub(r"thank\s+you","",t,flags=re.IGNORECASE)   # standard closing, not 1st/2nd person
        for m in set(w.lower() for w in re.findall(FORBID,t,flags=re.IGNORECASE)):
            voice.append((i,m))
        if re.search(r"(?<![\w/])I(?![\w/])",t):   # standalone capital I (skip I/O, IBM etc.)
            voice.append((i,"I"))
print("VOICE violations (1st/2nd person):",len(voice))
for v in sorted(set(voice))[:40]: print("   slide",v[0],"->",v[1])
