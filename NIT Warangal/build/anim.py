"""PowerPoint per-element / per-paragraph entrance animations for python-pptx slides.
Each item reveals on click, in order: whole shapes (fade/fly-in) or individual paragraphs of a
text shape (fade, line-by-line). Validated to open in Keynote/PowerPoint."""
from pptx.oxml.ns import qn
from pptx.oxml import parse_xml

P="http://schemas.openxmlformats.org/presentationml/2006/main"
A="http://schemas.openxmlformats.org/drawingml/2006/main"

def _tgt(spid, para):
    if para is None:
        return f'<p:tgtEl><p:spTgt spid="{spid}"/></p:tgtEl>'
    return f'<p:tgtEl><p:spTgt spid="{spid}"><p:txEl><p:pRg st="{para}" end="{para}"/></p:txEl></p:spTgt></p:tgtEl>'

def _click_par(cid, spid, effect, para):
    def nid():
        v=cid[0]; cid[0]+=1; return v
    outer=nid(); mid=nid(); eff=nid(); setc=nid(); anim1=nid()
    preset={"fade":10,"flyup":2,"flyleft":2,"flyright":2}[effect]
    tgt=_tgt(spid,para)
    motion=""
    if effect in ("flyup","flyleft","flyright"):
        anim2=nid()
        attr="ppt_y" if effect=="flyup" else "ppt_x"
        frm="0.3" if effect=="flyup" else ("-0.3" if effect=="flyleft" else "0.3")
        motion=f'''
          <p:anim calcmode="lin" valueType="num">
           <p:cBhvr additive="base">
            <p:cTn id="{anim2}" dur="450" fill="hold"/>{tgt}
            <p:attrNameLst><p:attrName>{attr}</p:attrName></p:attrNameLst>
           </p:cBhvr>
           <p:tavLst>
            <p:tav tm="0"><p:val><p:strVal val="{frm}"/></p:val></p:tav>
            <p:tav tm="100000"><p:val><p:strVal val="0"/></p:val></p:tav>
           </p:tavLst>
          </p:anim>'''
    return f'''
    <p:par>
     <p:cTn id="{outer}" fill="hold">
      <p:stCondLst><p:cond delay="indefinite"/></p:stCondLst>
      <p:childTnLst>
       <p:par>
        <p:cTn id="{mid}" fill="hold">
         <p:stCondLst><p:cond delay="0"/></p:stCondLst>
         <p:childTnLst>
          <p:par>
           <p:cTn id="{eff}" presetID="{preset}" presetClass="entr" presetSubtype="0" fill="hold" grpId="0" nodeType="clickEffect">
            <p:stCondLst><p:cond delay="0"/></p:stCondLst>
            <p:childTnLst>
             <p:set>
              <p:cBhvr>
               <p:cTn id="{setc}" dur="1" fill="hold"><p:stCondLst><p:cond delay="0"/></p:stCondLst></p:cTn>
               {tgt}
               <p:attrNameLst><p:attrName>style.visibility</p:attrName></p:attrNameLst>
              </p:cBhvr>
              <p:to><p:strVal val="visible"/></p:to>
             </p:set>
             <p:animEffect transition="in" filter="fade">
              <p:cBhvr><p:cTn id="{anim1}" dur="450"/>{tgt}</p:cBhvr>
             </p:animEffect>{motion}
            </p:childTnLst>
           </p:cTn>
          </p:par>
         </p:childTnLst>
        </p:cTn>
       </p:par>
      </p:childTnLst>
     </p:cTn>
    </p:par>'''

def add_entrance_anims(slide, items):
    """items: list of (shape, effect, para) — para is int paragraph index or None for whole shape."""
    items=[it for it in items if it[0] is not None]
    if not items: return
    cid=[3]
    clicks="".join(_click_par(cid, sh.shape_id, eff, para) for sh,eff,para in items)
    # one bldP per shape: build='p' if any paragraph-level item for it, else allAtOnce
    byshape={}
    for sh,eff,para in items:
        byshape.setdefault(sh.shape_id, False)
        if para is not None: byshape[sh.shape_id]=True
    bld="".join(
        (f'<p:bldP spid="{sid}" grpId="0" build="p" animBg="1"/>' if bp
         else f'<p:bldP spid="{sid}" grpId="0" animBg="1"/>')
        for sid,bp in byshape.items())
    timing=f'''<p:timing xmlns:p="{P}" xmlns:a="{A}">
     <p:tnLst>
      <p:par>
       <p:cTn id="1" dur="indefinite" restart="never" nodeType="tmRoot">
        <p:childTnLst>
         <p:seq concurrent="1" nextAc="seek">
          <p:cTn id="2" dur="indefinite" nodeType="mainSeq">
           <p:childTnLst>{clicks}</p:childTnLst>
          </p:cTn>
          <p:prevCondLst><p:cond evt="onPrev" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:prevCondLst>
          <p:nextCondLst><p:cond evt="onNext" delay="0"><p:tgtEl><p:sldTgt/></p:tgtEl></p:cond></p:nextCondLst>
         </p:seq>
        </p:childTnLst>
       </p:cTn>
      </p:par>
     </p:tnLst>
     <p:bldLst>{bld}</p:bldLst>
    </p:timing>'''
    el=parse_xml(timing)
    sld=slide.element
    ext=sld.find(qn('p:extLst'))
    if ext is not None: ext.addprevious(el)
    else: sld.append(el)

def add_transition(slide, kind="fade", spd="med"):
    """Add a slide-level page transition (fade between slides). Placed before <p:timing>."""
    sld=slide.element
    for t in sld.findall(qn('p:transition')): sld.remove(t)
    el=parse_xml(f'<p:transition xmlns:p="{P}" spd="{spd}"><p:{kind}/></p:transition>')
    timing=sld.find(qn('p:timing')); cmo=sld.find(qn('p:clrMapOvr'))
    if timing is not None: timing.addprevious(el)
    elif cmo is not None: cmo.addnext(el)
    else: sld.find(qn('p:cSld')).addnext(el)

if __name__=="__main__":
    from pptx import Presentation
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    from pptx.enum.shapes import MSO_SHAPE
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    s=prs.slides.add_slide(prs.slide_layouts[6])
    # a text box with 4 paragraphs (by-paragraph fade) + 2 shapes (fly)
    tb=s.shapes.add_textbox(Inches(0.6),Inches(0.6),Inches(7),Inches(4)); tf=tb.text_frame; tf.word_wrap=True
    for i in range(4):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); r=p.add_run(); r.text=f"Line {i+1}: verbose teaching sentence."; r.font.size=Pt(18)
    items=[(tb,"fade",i) for i in range(4)]
    for i in range(2):
        sh=s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE,Inches(8.5),Inches(0.6+i*2),Inches(3),Inches(1.5))
        sh.fill.solid(); sh.fill.fore_color.rgb=RGBColor(0x0F,0x62,0xFE); sh.text_frame.text=f"card {i+1}"
        items.append((sh,"flyup" if i==0 else "flyleft",None))
    add_entrance_anims(s,items)
    out="/private/tmp/claude-501/-Users-sentry-Work-Lectures/f0d8d36b-5746-4aa8-b83a-e2ca3d5eee0d/scratchpad/anim_test2.pptx"
    prs.save(out)
    Presentation(out); print("OK saved+reloaded para-level anims")
