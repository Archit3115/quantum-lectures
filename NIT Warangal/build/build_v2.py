#!/usr/bin/env python3
"""Build the verbose, animated Session-1 deck (v2) on a COPY of the IBM Qiskit template.
Data-driven: skeleton.json (layout/diagrams/numbers) + content.json (verbose read-aloud prose).
Adds per-paragraph / per-element entrance animations and Qiskit decorative assets to fill whitespace."""
import json, os
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn
from anim import add_entrance_anims, add_transition

ROOT="/Users/sentry/Work/Lectures"
TEMPLATE=os.path.join(ROOT,"Templates - IBM Qiskit Advocate","Qiskit_Foundations_QAP.pptx")
GEN=os.path.join(ROOT,"NIT Warangal","assets_generated")
ASSETS=os.path.join(ROOT,"assets")
TA=os.path.join(ROOT,"NIT Warangal","template_assets")
BUILD=os.path.dirname(__file__)
OUT=os.path.join(ROOT,"NIT Warangal","session1_quantum_clustering_qiskit.pptx")
def g(n): return os.path.join(GEN,n)
def a(*n): return os.path.join(ASSETS,*n)
def ta(n): return os.path.join(TA,n)

BLUE=RGBColor(0x0F,0x62,0xFE); PURPLE=RGBColor(0xA5,0x6E,0xFF); TEAL=RGBColor(0x00,0x9D,0x9A)
DEEP=RGBColor(0x00,0x3A,0x6D); MAGENTA=RGBColor(0x9F,0x18,0x53); RED=RGBColor(0xFA,0x4D,0x56)
INK=RGBColor(0x16,0x16,0x16); GRAY=RGBColor(0x6F,0x6F,0x6F); GRAY4=RGBColor(0x44,0x44,0x44)
LGRAY=RGBColor(0xC6,0xC6,0xC6); PANEL=RGBColor(0xF4,0xF4,0xF4); WHITE=RGBColor(0xFF,0xFF,0xFF)
CY=RGBColor(0x33,0xB1,0xFF); INDIGO=RGBColor(0x1A,0x0A,0x3E)
FONT="IBM Plex Sans"; MONO="IBM Plex Mono"
ACC={"blue":BLUE,"purple":PURPLE,"teal":TEAL,"deep":DEEP,"magenta":MAGENTA,"ink":INK,"red":RED}

SKEL=json.load(open(os.path.join(BUILD,"skeleton.json")))
CONTENT=json.load(open(os.path.join(BUILD,"content.json")))
def C(sid): return CONTENT.get(sid,{"body":[],"cards":[]})

# ---- strict margin grid (every slide obeys these) ----
MX=IN(0.55); RX=IN(12.78); FW=RX-MX          # left margin, right edge, full content width
CT=IN(1.66); CB=IN(7.02)                       # content top / bottom
GAP=IN(0.4); LW=IN(6.0); RCX=MX+LW+GAP; RCW=RX-RCX   # two-column: left 6.0in, right col
def strip_ph(slide,keep=()):
    for ph in list(slide.placeholders):
        if ph.placeholder_format.idx in keep: continue
        ph._element.getparent().remove(ph._element)

# ---------------- primitives ----------------
def L(prs,name):
    for lay in prs.slide_masters[0].slide_layouts:
        if lay.name==name: return lay
    raise KeyError(name)
def set_bg(slide,hexrgb):
    cSld=slide.element.find(qn('p:cSld')); old=cSld.find(qn('p:bg'))
    if old is not None: cSld.remove(old)
    bg=cSld.makeelement(qn('p:bg'),{}); bgPr=cSld.makeelement(qn('p:bgPr'),{})
    fill=cSld.makeelement(qn('a:solidFill'),{}); clr=cSld.makeelement(qn('a:srgbClr'),{'val':hexrgb})
    fill.append(clr); bgPr.append(fill); bgPr.append(cSld.makeelement(qn('a:effectLst'),{}))
    bg.append(bgPr); cSld.insert(0,bg)
def txt(slide,x,y,w,h,anchor=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    if anchor: tf.vertical_anchor=anchor
    return tb,tf
def _bold_split(text,bolds):
    marks=[]
    for b in (bolds or []):
        i=text.find(b)
        if i>=0: marks.append((i,i+len(b)))
    marks.sort(); out=[]; pos=0
    for s,e in marks:
        if s<pos: continue
        if s>pos: out.append((text[pos:s],False))
        out.append((text[s:e],True)); pos=e
    if pos<len(text): out.append((text[pos:],False))
    return out or [(text,False)]
def para(tf,text,size=12,color=INK,bold=False,italic=False,font=FONT,align=PP_ALIGN.LEFT,
         after=4,before=0,line=1.12,first=False,bolds=None,bullet=False,bcolor=None):
    p=tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment=align
    if after is not None: p.space_after=Pt(after)
    if before: p.space_before=Pt(before)
    if line: p.line_spacing=line
    if bullet:
        r=p.add_run(); r.text="•  "; r.font.size=Pt(size); r.font.bold=True
        r.font.color.rgb=bcolor or BLUE; r.font.name=font
    for chunk,isb in _bold_split(text,bolds):
        r=p.add_run(); r.text=chunk; f=r.font; f.size=Pt(size); f.name=font
        f.italic=italic; f.bold=(bold or isb); f.color.rgb=color
    return p
def rect(slide,x,y,w,h,fill,line=None,lw=1.0,round=False,rad=None):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,x,y,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    if round and rad is not None:
        try: shp.adjustments[0]=rad
        except Exception: pass
    return shp
def shp_text(shp,lines,anchor=MSO_ANCHOR.MIDDLE):
    tf=shp.text_frame; tf.word_wrap=True
    tf.margin_left=IN(0.08);tf.margin_right=IN(0.08);tf.margin_top=IN(0.04);tf.margin_bottom=IN(0.04)
    tf.vertical_anchor=anchor
    for i,(t,s,c,b) in enumerate(lines):
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=PP_ALIGN.CENTER
        r=p.add_run(); r.text=t; f=r.font; f.size=Pt(s); f.bold=b; f.color.rgb=c; f.name=FONT
    return shp
def pic(slide,path,x,y,w,h=None,valign="mid",halign="ctr"):
    p=slide.shapes.add_picture(path,x,y,width=w)
    if h is not None and p.height>h:
        r=h/p.height; p.width=int(p.width*r); p.height=h
    if halign=="ctr": p.left=int(x+(w-p.width)/2)
    if h is not None:
        p.top=int(y+(h-p.height)/2) if valign=="mid" else y
    return p
def pic_alpha(shape,pct):
    """set picture opacity (pct 0-100) via a:alphaModFix"""
    blip=shape._element.find('.//'+qn('a:blip'))
    if blip is None: return
    amf=blip.makeelement(qn('a:alphaModFix'),{'amt':str(int(pct*1000))}); blip.append(amf)
def arrow_r(slide,x,y,w,h,color=GRAY):
    s=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h); s.fill.solid()
    s.fill.fore_color.rgb=color; s.line.fill.background(); s.shadow.inherit=False
    try: s.adjustments[0]=0.5; s.adjustments[1]=0.6
    except Exception: pass
    return s

PAGENO=[0]
def pageno(slide):
    PAGENO[0]+=1
    tb,tf=txt(slide,RX-IN(0.7),IN(7.12),IN(0.7),IN(0.28))
    para(tf,f"{PAGENO[0]:02d}",9,GRAY,bold=True,align=PP_ALIGN.RIGHT,first=True,after=0)
    try:
        p=slide.shapes.add_picture(ta("qiskit_sphere.png"),MX,IN(7.06),height=IN(0.26)); pic_alpha(p,55)
    except Exception: pass

def header(slide,kicker,title,kcolor=BLUE,tsize=24):
    tb,tf=txt(slide,MX,IN(0.36),FW,IN(0.3))
    para(tf,(kicker or "").upper(),11,kcolor,bold=True,first=True,after=0)
    tb2,tf2=txt(slide,MX,IN(0.64),FW,IN(0.86))
    para(tf2,title,tsize,INK,bold=True,first=True,after=0,line=1.0)
    rect(slide,MX,IN(1.5),IN(1.15),IN(0.05),kcolor)
    return CT

def body_box(slide,x,y,w,h,blocks,anim,base=11.5,anchor=MSO_ANCHOR.TOP):
    """render body blocks as one textbox, one paragraph each; queue per-paragraph fade anims."""
    tb,tf=txt(slide,x,y,w,h,anchor=anchor)
    n=len(blocks)
    for i,blk in enumerate(blocks):
        t=blk.get("t","p"); x_=blk.get("x",""); bd=blk.get("b",[])
        first=(i==0)
        if t=="head": para(tf,x_,base+1.5,INK,bold=True,first=first,after=2.5,before=(0 if first else 5),bolds=bd)
        elif t=="lead": para(tf,x_,base+1,INK,first=first,after=3.5,line=1.14,bolds=bd)
        elif t=="b": para(tf,x_,base,GRAY4,first=first,after=3,line=1.12,bolds=bd,bullet=True)
        elif t=="note": para(tf,x_,base-1,GRAY,italic=True,first=first,after=2.5,line=1.1,bolds=bd)
        else: para(tf,x_,base,GRAY4,first=first,after=3.5,line=1.14,bolds=bd)
    for i in range(n): anim.append((tb,"fade",i))
    return tb

def cards_row(slide,x,y,w,h,titles_accents,lines_list,anim,tsize=13,bsize=10.3,gap=0.2):
    n=len(titles_accents); cw=int((w-IN(gap)*(n-1))/n)
    for i,(ttl,acc) in enumerate(titles_accents):
        cx=x+i*(cw+IN(gap)); col=ACC.get(acc,BLUE)
        sh=rect(slide,cx,y,cw,h,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.03)
        rect(slide,cx,y,IN(0.07),h,col)
        tb,tf=txt(slide,cx+IN(0.22),y+IN(0.14),cw-IN(0.36),h-IN(0.26))
        para(tf,ttl,tsize,INK,bold=True,first=True,after=4,line=1.02)
        lines=lines_list[i] if i<len(lines_list) else []
        for ln in lines: para(tf,ln,bsize,GRAY4,after=3,line=1.1)
        anim.append((sh,"flyup" if i%2==0 else "fade",None))
    return cw

def stat_tiles(slide,x,y,stats,anim,cols=4,tw=2.92,th=1.7,gap=0.16,nsize=25):
    for i,st in enumerate(stats):
        r,c=divmod(i,cols); cx=x+c*IN(tw+gap); cy=y+r*IN(th+0.18)
        col=ACC.get(st.get("accent","blue"),BLUE)
        sh=rect(slide,cx,cy,IN(tw),IN(th),PANEL,line=LGRAY,lw=0.75,round=True,rad=0.04)
        rect(slide,cx,cy,IN(tw),IN(0.09),col)
        tb,tf=txt(slide,cx+IN(0.16),cy+IN(0.15),IN(tw-0.3),IN(th-0.24),anchor=MSO_ANCHOR.TOP)
        para(tf,st["num"],nsize,col,bold=True,first=True,after=1)
        para(tf,st["label"],9.5,GRAY4,after=0,line=1.05)
        anim.append((sh,"fade",None))

def flowH(slide,x,y,w,h,steps,color,anim,fs=11):
    n=len(steps); ag=IN(0.24); sw=int((w-ag*(n-1))/n)
    for i,s in enumerate(steps):
        sx=x+i*(sw+ag); b=rect(slide,sx,y,sw,h,color,round=True,rad=0.07)
        shp_text(b,[(s,fs,WHITE,True)]); anim.append((b,"flyleft",None))
        if i<n-1: arrow_r(slide,sx+sw+IN(0.02),y+h//2-IN(0.09),ag-IN(0.04),IN(0.18),GRAY)

def flowV(slide,x,y,w,h,steps,color,anim,fs=11):
    n=len(steps); vg=IN(0.16); sh_=int((h-vg*(n-1))/n)
    for i,s in enumerate(steps):
        sy=y+i*(sh_+vg); b=rect(slide,x,sy,w,sh_,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.05)
        rect(slide,x,sy,IN(0.07),sh_,color)
        tb,tf=txt(slide,x+IN(0.2),sy,w-IN(0.32),sh_,anchor=MSO_ANCHOR.MIDDLE)
        para(tf,s,fs,INK,bold=True,first=True,after=0,line=1.05)
        anim.append((b,"fade",None))

def build_table(slide,x,y,w,h,data,header_fill=BLUE,fs=10.5):
    rows=len(data);cols=len(data[0])
    gt=slide.shapes.add_table(rows,cols,x,y,w,h); tbl=gt.table
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            cel=tbl.cell(i,j); cel.margin_left=IN(0.1);cel.margin_right=IN(0.08)
            cel.margin_top=IN(0.03);cel.margin_bottom=IN(0.03);cel.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf=cel.text_frame; tf.word_wrap=True; p=tf.paragraphs[0]
            r=p.add_run(); r.text=str(val); f=r.font; f.size=Pt(fs); f.name=FONT
            if i==0: cel.fill.solid();cel.fill.fore_color.rgb=header_fill;f.color.rgb=WHITE;f.bold=True
            else:
                cel.fill.solid();cel.fill.fore_color.rgb=WHITE if i%2 else PANEL
                f.color.rgb=INK if j==0 else GRAY4; f.bold=(j==0)
    return gt

DEFBG=RGBColor(0xED,0xF1,0xFA)
def defbox(slide,x,y,w,h,defs,anim=None,accent=DEEP,cols=1):
    """Key-terms box: header bar + plain-language definitions. defs=[{term,def}]."""
    if not defs: return
    sh=rect(slide,x,y,w,h,DEFBG,line=LGRAY,lw=0.75,round=True,rad=0.03)
    rect(slide,x,y,w,IN(0.32),accent,round=False)
    tb,tf=txt(slide,x+IN(0.18),y+IN(0.03),w-IN(0.3),IN(0.28))
    para(tf,"KEY TERMS",9.5,WHITE,bold=True,first=True,after=0)
    if cols==2 and len(defs)>2:
        half=(len(defs)+1)//2; cw=(w-IN(0.6))/2
        for ci,chunk in enumerate((defs[:half],defs[half:])):
            tb2,tf2=txt(slide,x+IN(0.18)+ci*(cw+IN(0.2)),y+IN(0.4),cw,h-IN(0.46))
            for i,d in enumerate(chunk):
                para(tf2,f"{d['term']} — {d['def']}",8.8,GRAY4,first=(i==0),after=3,line=1.04,bolds=[d['term']])
    else:
        tb2,tf2=txt(slide,x+IN(0.18),y+IN(0.4),w-IN(0.34),h-IN(0.46))
        for i,d in enumerate(defs):
            para(tf2,f"{d['term']} — {d['def']}",9,GRAY4,first=(i==0),after=3.5,line=1.06,bolds=[d['term']])
    if anim is not None: anim.append((sh,"fade",None))

LOGOS={"turning":["ibm","google","dwave"],"finance":["jpmorgan","ibm","qcware"],
 "pharma":["algorithmiq"],"climate":["quantinuum","totalenergies"],
 "whynow2":["google","dwave"],"ligo1":["caltech"],"ligo2":["caltech"],"cern":["cern"],
 "iter":["iter"],"skadune":["skao","fermilab"],"future1":["ibm","google","quantinuum","psiquantum"],
 "nqm1":["iisc"],"eco":["qpiai","qnu","ibm"],"benchmark":["ibm"]}
def logo_row(slide,names,y,anim=None,h=IN(0.36),right_x=None):
    paths=[os.path.join(TA,"logos",n+".png") for n in names]
    paths=[p for p in paths if os.path.exists(p)]
    if not paths: return
    cx=(right_x if right_x is not None else RX) - IN(0.52)*len(paths)
    for p in paths:
        try:
            pc=slide.shapes.add_picture(p,cx,y,height=h); cx=pc.left+pc.width+IN(0.16)
            if anim is not None: anim.append((pc,"fade",None))
        except Exception: pass

def _defh(cont):  # adaptive Key-terms box height (2-column)
    nd=len(cont.get("defs",[]))
    if not nd: return 0
    rows=(nd+1)//2
    return min(IN(1.62),IN(0.44)+IN(0.40)*rows)
def cbot(cont):  # content bottom — shrink when a Key-terms box is present
    dh=_defh(cont)
    return CB-dh-IN(0.14) if dh else CB
def extras(slide,sp,cont,anim):
    if sp["id"] in LOGOS: logo_row(slide,LOGOS[sp["id"]],IN(1.36),anim)
    dh=_defh(cont)
    if dh: defbox(slide,MX,CB-dh,FW,dh,cont["defs"],anim,cols=2)

def decor_sphere(slide,x,y,h,alpha=16,purple=True):
    try:
        p=slide.shapes.add_picture(ta("qiskit_sphere_purple.png" if purple else "qiskit_sphere.png"),x,y,height=h)
        pic_alpha(p,alpha); return p
    except Exception: return None
def decor_graph(slide,x,y,w,alpha=100,blue=False):
    try:
        p=slide.shapes.add_picture(ta("graph_blue.png" if blue else "graph_purple.png"),x,y,width=w)
        if alpha<100: pic_alpha(p,alpha); return p
    except Exception: return None

# ================= per-kind renderers =================
def r_divider(slide,sp):
    col=ACC.get(sp.get("accent","blue"),BLUE); set_bg(slide,"%02X%02X%02X"%(col[0],col[1],col[2]))
    decor_sphere(slide,IN(8.7),IN(2.2),IN(6.2),alpha=22,purple=True)
    tb,tf=txt(slide,IN(0.8),IN(1.0),IN(9),IN(2.4)); para(tf,sp["num"],150,WHITE,bold=True,first=True,after=0)
    tb2,tf2=txt(slide,IN(0.9),IN(3.8),IN(11.4),IN(1.6)); para(tf2,sp["title"],44,WHITE,bold=True,first=True,after=0,line=1.0)
    rect(slide,IN(0.94),IN(5.15),IN(2.0),IN(0.06),WHITE)
    tb3,tf3=txt(slide,IN(0.94),IN(5.45),IN(10.5),IN(1.2)); para(tf3,sp.get("summary",""),18,WHITE,first=True,line=1.2)
    tb4,tf4=txt(slide,IN(0.94),IN(6.95),IN(11),IN(0.3))
    para(tf4,"QML-2026 · Session 1 · Quantum Unsupervised Learning",10.5,WHITE,first=True,after=0)

def r_statement(slide,sp,anim):
    dark = sp.get("accent") in ("ink","deep")
    if dark:
        col=ACC.get(sp["accent"],INK); set_bg(slide,"%02X%02X%02X"%(col[0],col[1],col[2])); tc=WHITE; sub=LGRAY
    else:
        set_bg(slide,"F4F4F4"); tc=INK; sub=GRAY4
    img=sp.get("img")
    if img and img.get("f"):
        pic(slide,g(img["f"]),RX-IN(5.3),IN(2.5),IN(5.3),IN(3.6),valign="top")
        tw=IN(6.6)
    else:
        decor_sphere(slide,RX-IN(3.6),IN(4.2),IN(4.6),alpha=14 if not dark else 26); tw=FW
    rect(slide,MX,IN(1.7),IN(1.6),IN(0.06),ACC.get(sp.get("accent","purple"),PURPLE) if not dark else CY)
    tbk,tfk=txt(slide,MX,IN(1.2),IN(9),IN(0.4)); para(tfk,(sp.get("kicker") or " ").upper(),11,CY if dark else BLUE,bold=True,first=True,after=0)
    tb,tf=txt(slide,MX,IN(1.95),tw,IN(2.0)); para(tf,sp["title"],32,tc,bold=True,first=True,after=0,line=1.05)
    bt=body_box(slide,MX,IN(3.85),tw,IN(3.1),C(sp["id"])["body"],anim,base=11.5)
    for r in bt.text_frame.paragraphs:
        for rr in r.runs:
            if rr.font.color and rr.font.color.type is not None: pass
    # recolor body text to sub for dark bg
    if dark:
        for p in bt.text_frame.paragraphs:
            for rr in p.runs:
                if rr.font.bold: rr.font.color.rgb=WHITE
                else: rr.font.color.rgb=sub

def r_textimg(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    img=sp.get("img") or {}; side=img.get("side","right"); f=img.get("f")
    cont=C(sp["id"]); defs=cont.get("defs",[]); body=cont["body"]
    if sp["id"] in LOGOS: logo_row(slide,LOGOS[sp["id"]],IN(1.36),anim)
    fv=sp.get("flowV"); cards=sp.get("cards"); dh=_defh(cont)
    if fv:  # flow diagram in right column, defs below it, body full height left
        fh=IN(3.5) if defs else min(IN(4.1),CB-top)
        flowV(slide,RCX,top,RCW,fh,sp["flowV"],ACC.get(sp.get("accent","teal"),TEAL),anim,fs=11)
        if defs: defbox(slide,RCX,top+fh+IN(0.16),RCW,CB-top-fh-IN(0.16),defs,anim,cols=2)
        body_box(slide,MX,top,LW,CB-top,body,anim,base=(10.5 if defs else 11.5))
        return
    if (f and side=="top") or cards:   # bottom defs strip
        cb=CB-dh-IN(0.14) if defs else CB
        if f and side=="top":
            pic(slide,g(f),MX,top,FW,IN(2.1),valign="top")
            body_box(slide,MX,top+IN(2.26),FW,cb-top-IN(2.26),body,anim,base=(10 if defs else 11))
        else:
            cl=cont.get("cards",[])
            cards_row(slide,RCX,top,RCW,cb-top-IN(0.05),[(c["title"],c["accent"]) for c in cards],cl,anim,tsize=13,bsize=11)
            body_box(slide,MX,top,LW,cb-top,body,anim,base=(10 if defs else 11.5))
        if defs: defbox(slide,MX,CB-dh,FW,dh,defs,anim,cols=2)
        return
    # side image OR plain text: defs sit in the RIGHT COLUMN, body keeps FULL height
    if side=="left" and f: imx,bodyx,bodyw=MX,RCX,RCW
    else: imx,bodyx,bodyw=RCX,MX,(LW if (f or defs) else FW)
    dy=top
    if f:
        ih=IN(2.3) if defs else min(IN(3.95),CB-top)
        pic(slide,g(f),imx,top,RCW,ih,valign="top"); dy=top+ih+IN(0.16)
    if defs: defbox(slide,imx,dy,RCW,CB-dy,defs,anim,cols=1)
    elif not f: decor_graph(slide,RX-IN(3.5),top+IN(0.4),IN(3.5),alpha=100,blue=True)
    body_box(slide,bodyx,top,bodyw,CB-top,body,anim,base=(10.5 if defs else (11 if f else 11.8)))

def r_cards(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    cont=C(sp["id"]); body=cont["body"]; cl=cont.get("cards",[])
    n=len(sp["cards"]); cy=top; cb=cbot(cont); hasdef=bool(cont.get("defs"))
    if body:                                   # slim intro line above cards
        intro=body[:1] if hasdef else body[:2]
        body_box(slide,MX,cy,FW,IN(0.6) if hasdef else IN(0.72),intro,anim,base=11.5); cy=cy+(IN(0.62) if hasdef else IN(0.8))
    ch=cb-cy
    cards_row(slide,MX,cy,FW,ch,[(c["title"],c["accent"]) for c in sp["cards"]],cl,anim,
              tsize=13 if n<=3 else 12, bsize=(9 if hasdef else 10.6) if n<=3 else 9.2)
    extras(slide,sp,cont,anim)

def r_stats(slide,sp,anim,two_by_two=False):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    stats=sp["stats"]; cont=C(sp["id"]); defs=cont.get("defs",[])
    if sp["id"] in LOGOS: logo_row(slide,LOGOS[sp["id"]],IN(1.36),anim)
    stat_tiles(slide,MX,top,stats,anim,cols=4,tw=2.92,th=1.5,nsize=24)
    ty=top+IN(1.72)
    if defs:                      # body left, Key-terms in the right column
        body_box(slide,MX,ty,IN(7.4),CB-ty,cont["body"],anim,base=9.5)
        defbox(slide,IN(8.15),ty,RX-IN(8.15),CB-ty,defs,anim,cols=1)
    else:
        body_box(slide,MX,ty,FW,CB-ty,cont["body"],anim,base=11.5)

def r_table(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    cont=C(sp["id"]); defs=cont.get("defs",[])
    tx0=MX+IN(5.7); tw0=RX-tx0
    build_table(slide,tx0,top+IN(0.05),tw0,IN(3.25),sp["table"],header_fill=PURPLE,fs=9.5)
    tb,tf=txt(slide,tx0,top+IN(3.4),tw0,IN(0.4))
    para(tf,"*logarithmic scaling assumes efficient QRAM — see the reality check.",9,GRAY,italic=True,first=True,after=0)
    if defs: defbox(slide,tx0,top+IN(3.85),tw0,CB-top-IN(3.85),defs,anim,cols=2)
    body_box(slide,MX,top,IN(5.5),CB-top,cont["body"],anim,base=(10.5 if defs else 11))

def r_text(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    cont=C(sp["id"]); body=cont["body"]; cb=cbot(cont)
    if len(body)>6:
        half=(len(body)+1)//2
        body_box(slide,MX,top,IN(5.83),cb-top,body[:half],anim,base=(10.5 if cont.get("defs") else 11.5))
        body_box(slide,RCX,top,RCW,cb-top,body[half:],anim,base=(10.5 if cont.get("defs") else 11.5))
    else:
        body_box(slide,MX,top,IN(7.9),cb-top,body,anim,base=(10.7 if cont.get("defs") else 12))
        if not cont.get("defs"): decor_graph(slide,RX-IN(3.5),top+IN(0.6),IN(3.5),alpha=100,blue=False)
    extras(slide,sp,cont,anim)

def r_notebook(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"],kcolor=PURPLE,tsize=22)
    imgs=sp.get("imgs",[]); lit=sp.get("lit",[])
    if len(imgs)>=2:
        rcx=IN(7.35); rcw=RX-rcx
        p1=pic(slide,g(imgs[0]),MX,top,IN(6.5),IN(5.15),valign="top"); anim.append((p1,"fade",None))
        p2=pic(slide,g(imgs[1]),rcx,top,rcw,IN(2.3),valign="top"); anim.append((p2,"fade",None))
        body_box(slide,rcx,top+IN(2.5),rcw,CB-top-IN(2.5),lit,anim,base=11)
    elif imgs:
        rcx=IN(8.3); rcw=RX-rcx
        p1=pic(slide,g(imgs[0]),MX,top,IN(7.55),IN(5.15),valign="top"); anim.append((p1,"fade",None))
        body_box(slide,rcx,top,rcw,CB-top,lit,anim,base=12)

# ---- bespoke front/back ----
def r_title(slide,sp,anim):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx==11:
            try: ph.insert_picture(a("illustration-refs","hardware-blue.jpg"))
            except Exception: pass
    scr=rect(slide,0,0,IN(7.6),IN(7.5),INK)
    sf=scr.fill._xPr.find(qn('a:solidFill')); clr=sf.find(qn('a:srgbClr'))
    clr.append(clr.makeelement(qn('a:alpha'),{'val':'42000'}))
    tb,tf=txt(slide,IN(0.7),IN(1.05),IN(6.7),IN(3.7))
    para(tf,"QUANTUM MACHINE LEARNING · QML-2026",13,CY,bold=True,first=True,after=10)
    para(tf,"Quantum Unsupervised Learning",40,WHITE,bold=True,after=0,line=1.0)
    para(tf,"(Clustering)",40,CY,bold=True,after=10,line=1.0)
    for b in C("title")["body"][:1]:
        para(tf,b["x"],14.5,LGRAY,after=0,line=1.18,bolds=b.get("b"))
    tb2,tf2=txt(slide,IN(0.72),IN(5.75),IN(11),IN(1.5))
    para(tf2,"Archit Srivastava",18,WHITE,bold=True,first=True,after=2)
    para(tf2,"Senior Manager, Data Engineering @ PUMA · Founder, AiQyaM · Expert Session 1",12,LGRAY,after=2)
    para(tf2,"FDP · E&ICT Academy, NIT Warangal × NIT Raipur    |    11 July 2026 · 09:30–11:30",11.5,LGRAY,after=0)
    try:
        p=slide.shapes.add_picture(ta("qiskit_wordmark.png"),IN(10.7),IN(6.7),height=IN(0.5)); pic_alpha(p,85)
    except Exception: pass

def r_agenda(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    decor_sphere(slide,IN(10.8),IN(0.2),IN(1.3),alpha=30)
    secs=[("01","Introduction",BLUE),("02","Classical Unsupervised Learning",TEAL),
          ("03","Quantum Unsupervised Learning",PURPLE),("04","Industry & Current Research",MAGENTA),
          ("05","The Indian Landscape",DEEP),("06","The Future",INK)]
    desc=[b["x"] for b in C("agenda")["body"] if b["t"]=="b"]
    cw=IN(5.98); ch=IN(1.5); gx=IN(0.3); gy=IN(0.2)
    for i,(n,t,c) in enumerate(secs):
        r,cc=divmod(i,2); x=MX+cc*(cw+gx); y=top+r*(ch+gy)
        sh=rect(slide,x,y,cw,ch,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.05); rect(slide,x,y,IN(0.9),ch,c)
        tbn,tfn=txt(slide,x+IN(0.05),y,IN(0.86),ch,anchor=MSO_ANCHOR.MIDDLE)
        para(tfn,n,28,WHITE,bold=True,first=True,align=PP_ALIGN.CENTER,after=0)
        tbc,tfc=txt(slide,x+IN(1.05),y+IN(0.16),cw-IN(1.2),ch-IN(0.3),anchor=MSO_ANCHOR.MIDDLE)
        para(tfc,t,14.5,INK,bold=True,first=True,after=3,line=1.0)
        para(tfc,desc[i] if i<len(desc) else "",9.8,GRAY4,after=0,line=1.06)
        anim.append((sh,"fade",None))

def r_aboutbio(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    try: pic(slide,a("founder-archit.jpg"),MX,top,IN(2.55),IN(2.55),valign="top")
    except Exception: pass
    tb,tf=txt(slide,MX,top+IN(2.7),IN(2.6),IN(1.7))
    para(tf,"Archit Srivastava",13.5,INK,bold=True,first=True,after=2)
    para(tf,"Senior Manager, Data Engineering — PUMA, Bengaluru",10,GRAY4,after=3,line=1.05)
    para(tf,"Google Scholar: NbPUdWMAAAAJ",9.5,BLUE,after=0)
    body_box(slide,IN(3.45),top,RX-IN(3.45),CB-top,C("aboutbio")["body"],anim,base=12)
    try:
        p=slide.shapes.add_picture(ta("qiskit_sphere_purple.png"),IN(11.4),top+IN(3.3),height=IN(1.2)); pic_alpha(p,26)
    except Exception: pass

def r_aboutresearch(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    body_box(slide,MX,top,IN(7.0),CB-top,C("aboutresearch")["body"],anim,base=11.5)
    # publications card
    pubs=["Quantum Computing and LIGO — IAC-20, 2020","Quantum Finance: An Overview — EasyChair 6071, 2021",
          "ASD Therapy: AI-integrated Robotic Approach — IOP JPCS 2161, 2022","Data Analysis & Automation — EasyChair, 2021",
          "Shape-Memory Multi-Band Antenna — IAC-20, 2020","Smart Material for Satellite Antenna — IAC-19, 2019"]
    sh=rect(slide,IN(7.85),top,RX-IN(7.85),CB-top,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.03)
    rect(slide,IN(7.85),top,IN(0.07),CB-top,PURPLE)
    tb,tf=txt(slide,IN(8.07),top+IN(0.18),RX-IN(8.27),CB-top-IN(0.3))
    para(tf,"Selected publications  (Scholar: NbPUdWMAAAAJ)",12,INK,bold=True,first=True,after=6,line=1.05)
    for i,p_ in enumerate(pubs):
        para(tf,p_,10,GRAY4,after=5,line=1.1,bullet=True,bcolor=PURPLE)
    anim.append((sh,"fade",None))

def r_roadmap(slide,sp,anim):
    set_bg(slide,"FFFFFF"); top=header(slide,sp.get("kicker"),sp["title"])
    flowH(slide,MX,top+IN(0.12),FW,IN(1.1),
          ["The idea\n(unsupervised)","Classical tools\n(k-means, spectral)","The quantum turn\n(swap test, kernels)","The honest picture\n(demo + caveats)","Impact\n(industry, India)"],
          BLUE,anim,fs=12)
    body_box(slide,MX,top+IN(1.5),FW,CB-top-IN(1.5),C("roadmap")["body"],anim,base=12.5)

def r_thanks(slide,sp,anim):
    set_bg(slide,"161616"); rect(slide,0,0,IN(0.28),IN(7.5),BLUE)
    decor_sphere(slide,IN(9.0),IN(3.6),IN(4.2),alpha=20)
    tb,tf=txt(slide,IN(1.0),IN(1.05),IN(11),IN(2.0))
    para(tf,"Thank you.",46,WHITE,bold=True,first=True,after=8)
    for b in C("thanks")["body"][:2]:
        para(tf,b["x"],15,LGRAY,after=4,line=1.2,bolds=b.get("b"))
    tb2,tf2=txt(slide,IN(1.0),IN(3.7),IN(7.6),IN(3.2))
    para(tf2,"Archit Srivastava",20,WHITE,bold=True,first=True,after=6)
    para(tf2,"Senior Manager, Data Engineering @ PUMA  ·  Founder, AiQyaM",12.5,LGRAY,after=10)
    for lab,val in [("Email","architsrivastava3115@gmail.com"),("Scholar","NbPUdWMAAAAJ  (scholar.google.com)"),
                    ("Session","Quantum Unsupervised Learning (Clustering) · QML-2026 · 11 Jul 2026"),
                    ("Notebooks","runs on any laptop, no quantum hardware")]:
        para(tf2,f"{lab}   {val}",13,WHITE,after=4,bolds=[lab])
    try:
        p=slide.shapes.add_picture(ta("qiskit_wordmark.png"),IN(9.3),IN(6.7),height=IN(0.5)); pic_alpha(p,80)
    except Exception: pass

# ================= main build =================
prs=Presentation(TEMPLATE)
sldIdLst=prs.slides._sldIdLst
for sid in list(sldIdLst):
    prs.part.drop_rel(sid.get(qn('r:id'))); sldIdLst.remove(sid)
BLANK=L(prs,"Blank slide"); COVER=L(prs,"Cover, imagery")

for sp in SKEL:
    kind=sp["kind"]; anim=[]
    if kind=="title":
        s=prs.slides.add_slide(COVER); strip_ph(s,keep=(11,)); r_title(s,sp,anim); PAGENO[0]+=1
    else:
        s=prs.slides.add_slide(BLANK); strip_ph(s)
        if kind=="divider": r_divider(s,sp); PAGENO[0]+=1
        elif kind=="statement": r_statement(s,sp,anim); pageno(s)
        elif kind=="textimg": r_textimg(s,sp,anim); pageno(s)
        elif kind=="cards": r_cards(s,sp,anim); pageno(s)
        elif kind in ("stats","statstext"): r_stats(s,sp,anim); pageno(s)
        elif kind=="table": r_table(s,sp,anim); pageno(s)
        elif kind=="text": r_text(s,sp,anim); pageno(s)
        elif kind=="notebook": r_notebook(s,sp,anim); pageno(s)
        elif kind=="agenda": r_agenda(s,sp,anim); pageno(s)
        elif kind=="aboutbio": r_aboutbio(s,sp,anim); pageno(s)
        elif kind=="aboutresearch": r_aboutresearch(s,sp,anim); pageno(s)
        elif kind=="roadmap": r_roadmap(s,sp,anim); pageno(s)
        elif kind=="thanks": r_thanks(s,sp,anim); PAGENO[0]+=1
        else: r_text(s,sp,anim); pageno(s)
    if anim: add_entrance_anims(s,anim)
    add_transition(s,"fade")

prs.save(OUT)
print("SAVED",OUT,"slides",len(prs.slides._sldIdLst))
