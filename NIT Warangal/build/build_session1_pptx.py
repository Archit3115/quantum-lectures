#!/usr/bin/env python3
"""Build the Session-1 deck on a COPY of the IBM Qiskit template.
Quantum Unsupervised Learning (Clustering) — QML-2026 FDP, Archit Srivastava.
Runs on: python-pptx. Diagrams pre-rendered by gen_diagrams.py into ../assets_generated/.
Design: 3 template layouts inherit the Qiskit master (Cover imagery / Section divider / Blank slide);
all content is placed manually for density + control. IBM Carbon palette, IBM Plex Sans."""
import os
from pptx import Presentation
from pptx.util import Inches as IN, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---------- paths ----------
ROOT="/Users/sentry/Work/Lectures"
TEMPLATE=os.path.join(ROOT,"Templates - IBM Qiskit Advocate","Qiskit_Foundations_QAP.pptx")
GEN=os.path.join(ROOT,"NIT Warangal","assets_generated")
ASSETS=os.path.join(ROOT,"assets")
OUT=os.path.join(ROOT,"NIT Warangal","session1_quantum_clustering_qiskit.pptx")
def g(n): return os.path.join(GEN,n)
def a(*n): return os.path.join(ASSETS,*n)

# ---------- IBM Carbon palette ----------
BLUE=RGBColor(0x0F,0x62,0xFE); PURPLE=RGBColor(0xA5,0x6E,0xFF); TEAL=RGBColor(0x00,0x9D,0x9A)
DEEP=RGBColor(0x00,0x3A,0x6D); MAGENTA=RGBColor(0x9F,0x18,0x53); RED=RGBColor(0xFA,0x4D,0x56)
INK=RGBColor(0x16,0x16,0x16); GRAY=RGBColor(0x6F,0x6F,0x6F); GRAY4=RGBColor(0x52,0x52,0x52)
LGRAY=RGBColor(0xC6,0xC6,0xC6); PANEL=RGBColor(0xF4,0xF4,0xF4); PANEL2=RGBColor(0xE8,0xF1,0xFF)
WHITE=RGBColor(0xFF,0xFF,0xFF); CY=RGBColor(0x33,0xB1,0xFF); GRN=RGBColor(0x42,0xBE,0x65)
AMBER=RGBColor(0xB2,0x86,0x00)
QYGREEN=RGBColor(0x01,0xA9,0x82); INDIGO=RGBColor(0x1A,0x0A,0x3E)
FONT="IBM Plex Sans"; MONO="IBM Plex Mono"

# ---------- primitives ----------
def L(prs,name):
    for lay in prs.slide_masters[0].slide_layouts:
        if lay.name==name: return lay
    raise KeyError(name)

def txt(slide,x,y,w,h,anchor=None):
    tb=slide.shapes.add_textbox(x,y,w,h); tf=tb.text_frame; tf.word_wrap=True
    tf.margin_left=0;tf.margin_right=0;tf.margin_top=0;tf.margin_bottom=0
    if anchor: tf.vertical_anchor=anchor
    return tb,tf

def para(tf,text,size=12,color=INK,bold=False,italic=False,font=FONT,align=PP_ALIGN.LEFT,
         after=4,before=0,line=None,first=False):
    p=tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment=align
    if after is not None: p.space_after=Pt(after)
    if before: p.space_before=Pt(before)
    if line: p.line_spacing=line
    r=p.add_run(); r.text=text; f=r.font
    f.size=Pt(size); f.bold=bold; f.italic=italic; f.name=font; f.color.rgb=color
    return p

def rich(tf,segs,size=12,align=PP_ALIGN.LEFT,after=4,line=None,first=False,bullet=False):
    """segs: list of (text,color,bold[,font])"""
    p=tf.paragraphs[0] if (first and not tf.paragraphs[0].runs) else tf.add_paragraph()
    p.alignment=align
    if after is not None: p.space_after=Pt(after)
    if line: p.line_spacing=line
    if bullet:
        r=p.add_run(); r.text="•  "; r.font.size=Pt(size); r.font.color.rgb=BLUE; r.font.bold=True; r.font.name=FONT
    for s in segs:
        r=p.add_run(); r.text=s[0]; f=r.font; f.size=Pt(size); f.color.rgb=s[1]; f.bold=s[2]
        f.name=s[3] if len(s)>3 else FONT
    return p

def rect(slide,x,y,w,h,fill,line=None,lw=1.0,round=False,rad=None):
    shp=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if round else MSO_SHAPE.RECTANGLE,x,y,w,h)
    if fill is None: shp.fill.background()
    else: shp.fill.solid(); shp.fill.fore_color.rgb=fill
    if line is None: shp.line.fill.background()
    else: shp.line.color.rgb=line; shp.line.width=Pt(lw)
    shp.shadow.inherit=False
    if round and rad is not None:
        try: shp.adjustments[0]=rad
        except Exception: pass
    return shp

def alpha(shape,opacity_pct):
    """set fill opacity (100=opaque)"""
    sf=shape.fill._xPr.find(qn('a:solidFill')); clr=sf.find(qn('a:srgbClr'))
    e=clr.makeelement(qn('a:alpha'),{'val':str(int(opacity_pct*1000))}); clr.append(e)

def shp_text(shp,lines,size=12,color=WHITE,bold=True,align=PP_ALIGN.CENTER,anchor=MSO_ANCHOR.MIDDLE):
    tf=shp.text_frame; tf.word_wrap=True
    tf.margin_left=IN(0.06);tf.margin_right=IN(0.06);tf.margin_top=IN(0.03);tf.margin_bottom=IN(0.03)
    tf.vertical_anchor=anchor
    if isinstance(lines,str): lines=[(lines,size,color,bold)]
    for i,ln in enumerate(lines):
        t,s,c,b = ln if isinstance(ln,tuple) else (ln,size,color,bold)
        p=tf.paragraphs[0] if i==0 else tf.add_paragraph(); p.alignment=align
        r=p.add_run(); r.text=t; f=r.font; f.size=Pt(s); f.bold=b; f.color.rgb=c; f.name=FONT
    return shp

def pic(slide,path,x,y,w,h=None,valign="mid",halign="ctr"):
    p=slide.shapes.add_picture(path,x,y,width=w)
    if h is not None and p.height>h:
        r=h/p.height; p.width=int(p.width*r); p.height=h
    if halign=="ctr": p.left=int(x+(w-p.width)/2)
    if h is not None:
        if valign=="mid": p.top=int(y+(h-p.height)/2)
        elif valign=="top": p.top=y
    return p

def arrow_r(slide,x,y,w,h,color=GRAY):
    s=slide.shapes.add_shape(MSO_SHAPE.RIGHT_ARROW,x,y,w,h)
    s.fill.solid();s.fill.fore_color.rgb=color;s.line.fill.background();s.shadow.inherit=False
    try:s.adjustments[0]=0.5;s.adjustments[1]=0.55
    except Exception:pass
    return s
def arrow_d(slide,x,y,w,h,color=GRAY):
    s=slide.shapes.add_shape(MSO_SHAPE.DOWN_ARROW,x,y,w,h)
    s.fill.solid();s.fill.fore_color.rgb=color;s.line.fill.background();s.shadow.inherit=False
    return s

# ---------- slide scaffolding ----------
PAGENO=[0]
def pageno(slide):
    PAGENO[0]+=1
    tb,tf=txt(slide,IN(12.35),IN(7.06),IN(0.7),IN(0.3))
    para(tf,f"{PAGENO[0]:02d}",9,GRAY,bold=True,align=PP_ALIGN.RIGHT,first=True,after=0)

def title_of(slide):
    for sh in slide.shapes:
        if sh.is_placeholder and sh.placeholder_format.idx==0: return sh
    return None

def set_title(slide,text,size=25,color=INK):
    t=title_of(slide)
    t.left=IN(0.55);t.top=IN(0.6);t.width=IN(12.2);t.height=IN(0.9)
    tf=t.text_frame;tf.clear();tf.word_wrap=True
    tf.margin_left=0;tf.margin_top=0;tf.margin_bottom=0
    p=tf.paragraphs[0];p.alignment=PP_ALIGN.LEFT
    r=p.add_run();r.text=text;f=r.font;f.size=Pt(size);f.bold=True;f.color.rgb=color;f.name=FONT
    return t

def set_bg(slide,hexrgb):
    """set a true slide background (behind all shapes, incl. placeholders)"""
    cSld=slide.element.find(qn('p:cSld'))
    old=cSld.find(qn('p:bg'))
    if old is not None: cSld.remove(old)
    bg=cSld.makeelement(qn('p:bg'),{})
    bgPr=cSld.makeelement(qn('p:bgPr'),{})
    fill=cSld.makeelement(qn('a:solidFill'),{})
    clr=cSld.makeelement(qn('a:srgbClr'),{'val':hexrgb}); fill.append(clr)
    bgPr.append(fill); bgPr.append(cSld.makeelement(qn('a:effectLst'),{}))
    bg.append(bgPr); cSld.insert(0,bg)

def header(prs,kicker,title,kcolor=BLUE,tsize=25):
    """new content slide; returns (slide, content_top_emu)"""
    s=prs.slides.add_slide(L(prs,"Blank slide"))
    set_bg(s,"FFFFFF")
    set_title(s,title,tsize)
    tb,tf=txt(s,IN(0.57),IN(0.34),IN(11),IN(0.3))
    para(tf,kicker.upper(),11,kcolor,bold=True,first=True,after=0)
    rect(s,IN(0.58),IN(1.5),IN(1.15),IN(0.05),kcolor)
    pageno(s)
    return s,IN(1.68)

def divider(prs,num,title,summary,color=BLUE,tcolor=WHITE):
    s=prs.slides.add_slide(L(prs,"Blank slide"))
    rect(s,0,0,prs.slide_width,prs.slide_height,color)              # full bleed
    tb,tf=txt(s,IN(0.75),IN(1.0),IN(9),IN(2.4))
    para(tf,num,150,WHITE if tcolor==WHITE else INK,bold=True,first=True,after=0)
    tb2,tf2=txt(s,IN(0.85),IN(3.9),IN(11.4),IN(1.6))
    para(tf2,title,44,tcolor,bold=True,first=True,after=0)
    rect(s,IN(0.9),IN(5.15),IN(2.0),IN(0.06),WHITE if tcolor==WHITE else INK)
    tb3,tf3=txt(s,IN(0.9),IN(5.45),IN(11),IN(1.2))
    para(tf3,summary,18,tcolor,first=True,line=1.15)
    # section tag
    tb4,tf4=txt(s,IN(0.9),IN(6.9),IN(11),IN(0.3))
    para(tf4,"QML-2026 · Session 1 · Quantum Unsupervised Learning",11,
         WHITE if tcolor==WHITE else INK,first=True,after=0)
    # opacity of tag
    alpha_para(tf4,70)
    PAGENO[0]+=1
    return s
def alpha_para(tf,pct):
    for p in tf.paragraphs:
        for r in p.runs:
            sf=r.font.color._xFill
            # skip; simple approach: leave as is
            pass

def card(slide,x,y,w,h,title,body,accent=BLUE,tsize=13,bsize=10.5,fill=PANEL,round=True):
    rect(slide,x,y,w,h,fill,line=LGRAY,lw=0.75,round=round,rad=0.04)
    rect(slide,x,y,IN(0.07),h,accent)
    tb,tf=txt(slide,x+IN(0.24),y+IN(0.14),w-IN(0.38),h-IN(0.26))
    para(tf,title,tsize,INK,bold=True,first=True,after=4)
    for ln in body:
        if isinstance(ln,tuple): rich(tf,[ln[0]] if False else ln,bsize)  # not used
        else: para(tf,ln,bsize,GRAY4,after=2.5,line=1.05)
    return

def stat(slide,x,y,w,h,number,label,color=BLUE,nsize=27,fill=PANEL):
    rect(slide,x,y,w,h,fill,line=LGRAY,lw=0.75,round=True,rad=0.05)
    rect(slide,x,y,w,IN(0.09),color)
    tb,tf=txt(slide,x+IN(0.16),y+IN(0.13),w-IN(0.3),h-IN(0.24),anchor=MSO_ANCHOR.TOP)
    para(tf,number,nsize,color,bold=True,first=True,after=1)
    para(tf,label,10,GRAY4,after=0,line=1.02)

def bullets(slide,x,y,w,h,items,size=12,color=GRAY4,gap=4,head=None,hcolor=INK):
    tb,tf=txt(slide,x,y,w,h)
    first=True
    if head:
        para(tf,head,size+2,hcolor,bold=True,first=True,after=6); first=False
    for it in items:
        segs=it if isinstance(it,list) else [(it,color,False)]
        rich(tf,segs,size,after=gap,bullet=True,line=1.05,first=first)
        first=False
    return tf

def flow(slide,x,y,w,h,steps,color=BLUE,tcol=WHITE,fs=12,agap=0.28):
    n=len(steps); ag=IN(agap); sw=int((w-ag*(n-1))/n)
    for i,s in enumerate(steps):
        sx=x+i*(sw+ag)
        b=rect(slide,sx,y,sw,h,color,round=True,rad=0.08)
        shp_text(b,s,fs,tcol,bold=True)
        if i<n-1: arrow_r(slide,sx+sw+IN(0.02),y+h//2-IN(0.09),ag-IN(0.04),IN(0.18),GRAY)

def build_table(slide,x,y,w,h,data,colw=None,header_fill=BLUE,fs=10.5):
    rows=len(data);cols=len(data[0])
    gtbl=slide.shapes.add_table(rows,cols,x,y,w,h);tbl=gtbl.table
    if colw:
        for j,cw in enumerate(colw): tbl.columns[j].width=cw
    for i,row in enumerate(data):
        for j,val in enumerate(row):
            c=tbl.cell(i,j);c.margin_left=IN(0.1);c.margin_right=IN(0.08)
            c.margin_top=IN(0.04);c.margin_bottom=IN(0.04);c.vertical_anchor=MSO_ANCHOR.MIDDLE
            tf=c.text_frame;tf.word_wrap=True;p=tf.paragraphs[0]
            r=p.add_run();r.text=str(val);f=r.font;f.size=Pt(fs);f.name=FONT
            if i==0:
                c.fill.solid();c.fill.fore_color.rgb=header_fill;f.color.rgb=WHITE;f.bold=True
            else:
                c.fill.solid();c.fill.fore_color.rgb=WHITE if i%2 else PANEL
                f.color.rgb=INK if j==0 else GRAY4; f.bold=(j==0)
    return tbl

# =====================================================================
#  BUILD
# =====================================================================
prs=Presentation(TEMPLATE)
# delete all example slides — drop the relationship (not just the sldId) so the
# orphaned slide parts are not re-serialized (avoids duplicate-partname collisions).
sldIdLst=prs.slides._sldIdLst
for sid in list(sldIdLst):
    prs.part.drop_rel(sid.get(qn('r:id')))
    sldIdLst.remove(sid)

SW,SH=prs.slide_width,prs.slide_height

# ---- 1. TITLE (Cover, imagery) ----
s=prs.slides.add_slide(L(prs,"Cover, imagery"))
for ph in s.placeholders:
    if ph.placeholder_format.idx==11:
        try: ph.insert_picture(a("illustration-refs","hardware-blue.jpg"))
        except Exception: pass
# scrim for legibility
scr=rect(s,0,0,IN(7.4),SH,INK); alpha(scr,42)
def cover_text(idx,lines):
    for ph in s.placeholders:
        if ph.placeholder_format.idx==idx:
            tf=ph.text_frame;tf.clear()
            for i,(t,sz,c,b) in enumerate(lines):
                p=tf.paragraphs[0] if i==0 else tf.add_paragraph()
                r=p.add_run();r.text=t;f=r.font;f.size=Pt(sz);f.color.rgb=c;f.bold=b;f.name=FONT
# title placeholder (idx 0)
tb,tf=txt(s,IN(0.7),IN(1.15),IN(6.7),IN(3.6))
para(tf,"QUANTUM MACHINE LEARNING · QML-2026",13,CY,bold=True,first=True,after=10)
para(tf,"Quantum Unsupervised Learning",40,WHITE,bold=True,after=0,line=1.0)
para(tf,"(Clustering)",40,CY,bold=True,after=10,line=1.0)
para(tf,"Building quantum clustering from first principles — no black boxes.",15,LGRAY,after=0,line=1.15)
tb2,tf2=txt(s,IN(0.72),IN(5.9),IN(11),IN(1.4))
para(tf2,"Archit Srivastava",18,WHITE,bold=True,first=True,after=2)
para(tf2,"Founder, AiQyaM · Expert Session 1",12.5,LGRAY,after=2)
para(tf2,"FDP · E&ICT Academy, NIT Warangal × NIT Raipur    |    11 July 2026 · 09:30–11:30",11.5,LGRAY,after=0)
PAGENO[0]+=1

# ---- 2. AGENDA ----
s,ct=header(prs,"Contents","What we'll build in the next two hours")
secs=[("01","Introduction","Why quantum ML now — the 2025 inflection, the market, and the data deluge",BLUE),
      ("02","Classical Unsupervised Learning","k-means, the landscape, spectral clustering, and where the cost lives",TEAL),
      ("03","Quantum Unsupervised Learning","Qubits, the swap test, quantum kernels, q-means, a live demo, the honest picture",PURPLE),
      ("04","Industry & Current Research","Finance, pharma, climate — and quantum ML inside LIGO, CERN, ITER, SKA, DUNE",MAGENTA),
      ("05","The Indian Landscape","National Quantum Mission, indigenous hardware, Viksit Bharat 2047",DEEP),
      ("06","The Future","Fault-tolerance, AI×Quantum, the talent gap — and your role in it",INK)]
cw=IN(5.95); ch=IN(1.55); gx=IN(0.3); gy=IN(0.22)
for i,(n,t,d,c) in enumerate(secs):
    r,cc=divmod(i,2)
    x=IN(0.57)+cc*(cw+gx); y=ct+r*(ch+gy)
    rect(s,x,y,cw,ch,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.05); rect(s,x,y,IN(0.9),ch,c)
    tbn,tfn=txt(s,x+IN(0.06),y+IN(0.18),IN(0.86),IN(1.2),anchor=MSO_ANCHOR.MIDDLE)
    para(tfn,n,30,WHITE,bold=True,first=True,align=PP_ALIGN.CENTER,after=0)
    tbc,tfc=txt(s,x+IN(1.06),y+IN(0.16),cw-IN(1.2),ch-IN(0.3))
    para(tfc,t,15,INK,bold=True,first=True,after=3,line=1.0)
    para(tfc,d,10.5,GRAY4,after=0,line=1.05)

# ---- 3. ABOUT ----
s,ct=header(prs,"About the speaker","Archit Srivastava")
try: pic(s,a("founder-archit.jpg"),IN(0.6),ct+IN(0.1),IN(2.5),IN(2.5),valign="top")
except Exception: pass
tb,tf=txt(s,IN(0.6),ct+IN(2.75),IN(2.5),IN(1.6))
para(tf,"Archit Srivastava",15,INK,bold=True,first=True,after=2)
para(tf,"Founder — AiQyaM (India's first quantum-hardware community, est. 2020)",10.5,GRAY4,after=3,line=1.05)
para(tf,"Google Scholar: NbPUdWMAAAAJ",10,BLUE,after=0)
cols=[("Quantum community",BLUE,["Founder, AiQyaM — India's first quantum-hardware community (2020)",
        "Founded CIRQuIT quantum-computing group at RV College of Engineering",
        "Runs cohorts that put students on real quantum hardware"]),
      ("Industry practice",TEAL,["Quantum & ML systems in industry — o9 Solutions",
        "Hewlett Packard Enterprise quantum proofs-of-concept",
        "\"I ship these methods, not just teach them\""]),
      ("Published research",PURPLE,["Quantum Finance — VQE portfolio optimization in BOTH Qiskit & PennyLane (EasyChair 6071, 2021)",
        "AI computer vision + Quantum Visual Tracking, run on real IBM hardware (IOP, J. Phys. Conf. 2161, 2022)",
        "Today's material is drawn from methods I've built and run"])]
x0=IN(3.45); cwd=IN(3.02); gp=IN(0.2)
for i,(h,c,items) in enumerate(cols):
    x=x0+i*(cwd+gp)
    card(s,x,ct+IN(0.1),cwd,IN(3.9),h,items,accent=c,tsize=13,bsize=10.2)
# Qyros wordmark bottom-right (height-fit so it never overflows)
try: s.shapes.add_picture(a("logo-wordmark-indigo.png"),IN(10.55),ct+IN(4.05),height=IN(0.5))
except Exception: pass

# ---- 4. ROADMAP ----
s,ct=header(prs,"How this session works","Start from what you already know — add one quantum idea at a time")
flow(s,IN(0.6),ct+IN(0.35),IN(12.2),IN(1.15),
     ["The idea\n(unsupervised)","Classical tools\n(k-means, spectral)","The quantum turn\n(swap test, kernels)","The honest picture\n(demo + caveats)","Impact\n(industry, India)"],
     color=BLUE,fs=12.5)
tb,tf=txt(s,IN(0.6),ct+IN(2.1),IN(12.2),IN(2.8))
para(tf,"The design principle",15,INK,bold=True,first=True,after=6)
rich(tf,[("Every quantum idea in this talk is introduced from a classical anchor you already teach. ",GRAY4,False),
         ("If you can teach k-means, you will leave able to teach the quantum version.",INK,True)],13,after=10,line=1.2)
rich(tf,[("The spine of the whole session, in one sentence:  ",GRAY4,False),
         ("clustering = representation + similarity + objective, and quantum mostly changes the similarity — and how fast we compute it.",INK,True)],13,after=10,line=1.2)
rich(tf,[("Two live notebooks (PennyLane, CPU-only) — ",GRAY4,False),("nothing here needs quantum hardware.",TEAL,True)],13,after=0,line=1.2)

# =====================================================================
# SECTION 1 — INTRODUCTION
# =====================================================================
divider(prs,"01","Introduction","Why should an Indian university invest attention in quantum machine learning — now?",BLUE)

# 1.1 Why now
s,ct=header(prs,"Section 1 · Introduction","2025: the year quantum advantage became measured fact")
pic(s,g("timeline.png"),IN(0.5),ct+IN(0.05),IN(12.4),IN(2.6),valign="top")
tb,tf=txt(s,IN(0.6),ct+IN(2.95),IN(12.2),IN(2.2))
rich(tf,[("Willow (Google, Nature Dec 2024): ",INK,True),("105 qubits, first below-threshold error correction — logical error shrinks as the code grows (Λ≈2.14).",GRAY4,False)],11.5,bullet=True,first=True,after=4,line=1.1)
rich(tf,[("Quantum Echoes (Google, Nature Oct 2025): ",INK,True),("first verifiable quantum advantage, ~13,000× a classical supercomputer, via out-of-time-order correlators.",GRAY4,False)],11.5,bullet=True,after=4,line=1.1)
rich(tf,[("D-Wave Advantage2 (Science Mar 2025): ",INK,True),("magnetic-materials simulation beyond classical reach (claim contested by tensor-network rebuttals — cite honestly).",GRAY4,False)],11.5,bullet=True,after=4,line=1.1)
rich(tf,[("For faculty: ",MAGENTA,True),("this is no longer a physics curiosity — it is a computational tool your students will be asked about in interviews next year.",INK,False)],11.5,bullet=True,after=0,line=1.1)

# 1.2 Market
s,ct=header(prs,"Section 1 · Introduction","The stakes: a commercial tipping point")
tiles=[("$1.3–2.7T","economic value by 2035 (McKinsey QTM 2026)",BLUE),
       ("$12.6B","quantum start-up investment in 2025 — 6.3× the 2024 figure",TEAL),
       (">$1B","quantum-computing revenue in 2025 → ~$4.4B by 2028",PURPLE),
       ("$450–850B","BCG's economic-value estimate by 2040",MAGENTA)]
tw=IN(2.92);th=IN(1.7);gx=IN(0.16)
for i,(n,l,c) in enumerate(tiles):
    stat(s,IN(0.57)+i*(tw+gx),ct+IN(0.1),tw,th,n,l,color=c,nsize=26)
tb,tf=txt(s,IN(0.6),ct+IN(2.05),IN(12.2),IN(2.6))
para(tf,"The line that lands in a faculty room — talent",14,INK,bold=True,first=True,after=6)
rich(tf,[("Industry projects ~250,000 quantum roles by 2030 and ~840,000 by 2035. ",INK,True),
         ("In 2025 there were roughly 3 open positions per qualified hire, and only about half of quantum roles could be filled.",GRAY4,False)],12.5,after=8,line=1.2,bullet=True)
rich(tf,[("That gap is the mandate. ",MAGENTA,True),
         ("The value of quantum is not in selling processors — it is in the algorithms, software and AI layers that capture the outcomes. That layer is a keyboard-and-brains problem — exactly India's strength.",GRAY4,False)],12.5,after=0,line=1.2,bullet=True)

# 1.3 Data deluge
s,ct=header(prs,"Section 1 · Introduction","Why unsupervised learning: the data is unlabelled")
pic(s,g("data_deluge.png"),IN(6.7),ct+IN(0.15),IN(6.1),IN(3.9),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.15),IN(5.9),IN(4.2),[
    [("~80% of real-world data is unlabelled",INK,True)],
    [("Labels are expensive, scarce, or impossible — clustering is how you make the first sense of a new dataset",GRAY4,False)],
    [("Customer segments · genomic subtypes · anomaly/fraud detection · document topics · image compression",GRAY4,False)],
    [("The curse of dimensionality: ",INK,True),("in high-D, nearest and farthest points become nearly equidistant — Euclidean distance loses its discriminative power",GRAY4,False)],
    [("Mega-science makes this extreme: ",MAGENTA,True),("SKA streams ~petabytes/second, the HL-LHC will archive exabytes/year, one RNA-seq sample spans 20,000+ genes",GRAY4,False)],
],size=12,gap=8)

# 1.4 Sorting coins
s,ct=header(prs,"Section 1 · Introduction","Unsupervised learning = sorting coins in the dark")
tb,tf=txt(s,IN(0.6),ct+IN(0.1),IN(6.0),IN(4.2))
para(tf,"The one-sentence definition",14,INK,bold=True,first=True,after=5)
para(tf,"Unsupervised learning finds structure in data that has no labels. You are handed a pile of points and asked: which of these belong together?",12.5,GRAY4,after=12,line=1.2)
para(tf,"The intuition anchor",14,INK,bold=True,after=5)
para(tf,"Tip a bag of mixed coins onto a table in the dark. You can't read the denominations (no labels), but you can feel size and weight (features). You sort them into piles by similarity. That sorting — no labels, only a notion of 'how alike' — is clustering.",12.5,GRAY4,after=0,line=1.2)
# ML map (native)
mx=IN(7.0)
rect(s,mx,ct+IN(0.2),IN(5.7),IN(0.7),INK,round=True,rad=0.08)
shp_text(s.shapes[-1],"Machine Learning",15,WHITE)
kids=[("Supervised","data + labels → predict\n(classification, regression)",TEAL),
      ("Unsupervised","data only → discover structure\n(clustering, dim-reduction)",BLUE),
      ("Reinforcement","learn from reward\n(control, games)",PURPLE)]
kw=IN(1.8);kh=IN(1.5);kg=IN(0.15)
for i,(t,d,c) in enumerate(kids):
    x=mx+i*(kw+kg)
    arrow_d(s,x+kw//2-IN(0.08),ct+IN(0.95),IN(0.16),IN(0.3),LGRAY)
    b=rect(s,x,ct+IN(1.35),kw,kh,c if t=="Unsupervised" else PANEL,round=True,rad=0.06)
    shp_text(b,[(t,12.5,WHITE if t=="Unsupervised" else INK,True),(d,9,WHITE if t=="Unsupervised" else GRAY4,False)],anchor=MSO_ANCHOR.MIDDLE)
tb,tf=txt(s,mx,ct+IN(3.05),IN(5.7),IN(1.0))
para(tf,"~80% of real data lives on the middle branch — no labels to learn from.",11,GRAY4,italic=True,first=True,after=0,line=1.1)

# 1.5 Three ingredients
s,ct=header(prs,"Section 1 · Introduction","Every clustering method is three choices")
ings=[("1 · Representation",BLUE,["Each object → a vector of features","A point in ℝᵈ","e.g. a customer as (age, spend, recency, …)"]),
      ("2 · Similarity / distance",MAGENTA,["How we measure \"alike\"","Euclidean distance, cosine, or a kernel","THIS is what quantum changes"]),
      ("3 · Objective / rule",TEAL,["What makes a grouping good","Minimise within-cluster spread","Maximise between-cluster separation"])]
cwd=IN(3.95);ch=IN(2.5)
for i,(h,c,items) in enumerate(ings):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=15,bsize=11.5)
tb,tf=txt(s,IN(0.6),ct+IN(2.85),IN(12.2),IN(1.4))
rich(tf,[("The thesis of the whole session:  ",INK,True),
         ("everything that follows — classical and quantum — is a different choice of these three. Quantum methods mostly change ingredient 2 (the similarity) and how fast we compute it. Hold that thought; it is the spine.",GRAY4,False)],13.5,first=True,after=0,line=1.25)

# =====================================================================
# SECTION 2 — CLASSICAL
# =====================================================================
divider(prs,"02","Classical Unsupervised Learning","The honest version of the tools you already teach — and where they hit a wall.",TEAL)

# 2.1 k-means Lloyd loop
s,ct=header(prs,"Section 2 · Classical","k-means: nearest-centroid reasoning")
flow(s,IN(0.6),ct+IN(0.1),IN(12.2),IN(1.05),
     ["Initialise\nk centroids\n(k-means++)","Assign\neach point →\nnearest centroid","Update\ncentroid →\nmean of its points","Repeat\nuntil assignments\nstop changing"],
     color=TEAL,fs=12)
tb,tf=txt(s,IN(0.6),ct+IN(1.55),IN(6.1),IN(2.9))
para(tf,"The objective (inertia / within-cluster sum of squares)",12.5,INK,bold=True,first=True,after=5)
para(tf,"J = Σ_clusters Σ_(x in cluster)  ‖x − μ_cluster‖²",14,BLUE,bold=True,font=MONO,after=10)
rich(tf,[("The assign step is n × k distance computations. ",MAGENTA,True),("Remember that line — it is exactly where quantum will push.",GRAY4,False)],12,after=0,line=1.2)
card(s,IN(6.9),ct+IN(1.55),IN(5.9),IN(2.9),"Strengths & the blind spot",[
    "Strengths: simple, fast, scales to millions of points",
    "Weaknesses: must fix k in advance; assumes round, equal-size, linearly-separable blobs",
    "Sensitive to initialisation and outliers",
    "Fails on non-convex shapes — interleaved crescents, concentric rings — because 'nearest centroid' is a straight-line idea"],
    accent=RED,bsize=11)

# 2.2 landscape
s,ct=header(prs,"Section 2 · Classical","Beyond k-means — the landscape, and the bridge to quantum")
land=[("Hierarchical / agglomerative",BLUE,["Build a tree (dendrogram) of merges","No need to fix k up front","O(n²)–O(n³)"]),
      ("DBSCAN",TEAL,["Density-based","Finds arbitrary shapes + labels outliers","Needs a density scale ε"]),
      ("Gaussian Mixture (GMM)",PURPLE,["Soft, probabilistic clusters","Expectation–Maximisation","Overlapping membership"]),
      ("Spectral clustering",MAGENTA,["Similarity graph → Laplacian eigenvectors → k-means in that space","Handles non-convex data","Heavy step = linear algebra on a similarity matrix"])]
cwd=IN(2.95);ch=IN(2.55)
for i,(h,c,items) in enumerate(land):
    card(s,IN(0.57)+i*(cwd+IN(0.15)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=12.5,bsize=10)
tb,tf=txt(s,IN(0.6),ct+IN(2.9),IN(12.2),IN(1.2))
rich(tf,[("Flag this: ",MAGENTA,True),("spectral clustering's heavy step is linear algebra on a similarity / kernel matrix — and a quantum computer is a native linear-algebra machine. That is the natural bridge to the quantum turn.",GRAY4,False)],13,first=True,after=0,line=1.2)

# 2.3 spectral bridge
s,ct=header(prs,"Section 2 · Classical","Spectral clustering — the bridge to quantum")
pic(s,g("separability.png"),IN(6.5),ct+IN(0.2),IN(6.3),IN(3.4),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.15),IN(5.7),IN(4.0),[
    [("Straight-line 'nearest centroid' thinking fails on curved data (left)",GRAY4,False)],
    [("Spectral clustering builds a similarity graph, embeds using the top eigenvectors of its graph Laplacian, then runs k-means in that embedded space",GRAY4,False)],
    [("Lift the data into a richer space and the clusters become separable (right)",GRAY4,False)],
    [("The heavy step is eigen-decomposition of an n×n similarity matrix",INK,True)],
    [("A quantum feature map produces a richer similarity for exactly this pipeline — same spectral method, new kernel",TEAL,True)],
],size=11.5,gap=8)

# 2.4 where cost lives
s,ct=header(prs,"Section 2 · Classical","Where the classical cost lives")
pic(s,g("complexity_ladder.png"),IN(0.5),ct+IN(0.1),IN(7.0),IN(3.9),valign="top")
tf=bullets(s,IN(7.7),ct+IN(0.3),IN(5.1),IN(4.0),[
    [("one Euclidean distance = O(d)",INK,True,MONO)],
    [("one k-means iteration = O(n·k·d)",INK,True,MONO)],
    [("build a kernel matrix = O(n²·d)",INK,True,MONO)],
    [("eigen-decompose it = O(n³)",INK,True,MONO)],
    [("Everything expensive is computing similarities between many high-dimensional vectors, and linear algebra on n×n matrices",GRAY4,False)],
    [("That is precisely the corner quantum targets — distance / inner-product estimation and high-dimensional feature spaces",MAGENTA,True)],
    [("Now we've earned the quantum part.",INK,True)],
],size=11.5,gap=7)

# =====================================================================
# SECTION 3 — QUANTUM
# =====================================================================
divider(prs,"03","Quantum Unsupervised Learning","Just enough quantum, from first principles — no hand-waving, no black boxes.",PURPLE)

# 3.0 QM as a tool
s=prs.slides.add_slide(L(prs,"Blank slide"))
rect(s,0,0,SW,SH,PANEL)
tb,tf=txt(s,IN(1.0),IN(2.2),IN(11.3),IN(3.2),anchor=MSO_ANCHOR.MIDDLE)
para(tf,"Quantum mechanics as a tool, not math.",34,INK,bold=True,first=True,after=14,line=1.05)
para(tf,"Quantum clustering does not execute classical arithmetic faster. It uses physical phenomena — superposition, interference, entanglement — to perform geometric evaluations that cannot be efficiently mapped onto classical von Neumann hardware.",16,GRAY4,after=0,line=1.3)
rect(s,IN(1.02),IN(2.0),IN(1.6),IN(0.06),PURPLE)
pageno(s)

# 3.1 qubit + bloch
s,ct=header(prs,"Section 3 · Quantum","A qubit is a coin that is every angle at once")
pic(s,g("bloch_sphere.png"),IN(8.2),ct+IN(0.0),IN(4.5),IN(4.2),valign="top")
tb,tf=txt(s,IN(0.6),ct+IN(0.1),IN(7.3),IN(4.2))
para(tf,"|ψ⟩ = α|0⟩ + β|1⟩,   with |α|² + |β|² = 1",15,BLUE,bold=True,font=MONO,first=True,after=10)
rich(tf,[("α, β are amplitudes (complex). ",INK,True),("|α|² is the probability of measuring 0.",GRAY4,False)],12.5,bullet=True,after=6,line=1.15)
rich(tf,[("Superposition: ",INK,True),("before measurement the qubit genuinely holds a weighted combination of both states.",GRAY4,False)],12.5,bullet=True,after=6,line=1.15)
rich(tf,[("Bloch sphere: ",INK,True),("|0⟩ is the north pole, |1⟩ the south — everything else is a direction. A qubit carries a continuous direction, not one bit.",GRAY4,False)],12.5,bullet=True,after=6,line=1.15)
rich(tf,[("This continuous, complex-valued state is the raw material quantum clustering works in.",MAGENTA,True)],12.5,bullet=True,after=0,line=1.15)

# 3.2 2^n + measurement
s,ct=header(prs,"Section 3 · Quantum","n qubits → 2ⁿ amplitudes — the resource, and the catch")
pic(s,g("amplitude_growth.png"),IN(0.5),ct+IN(0.1),IN(7.0),IN(3.9),valign="top")
tf=bullets(s,IN(7.7),ct+IN(0.3),IN(5.1),IN(4.0),[
    [("n qubits live in a space of 2ⁿ complex amplitudes",INK,True)],
    [("50 qubits ≈ 10¹⁵ amplitudes evolving together — this exponential state space is the resource",GRAY4,False)],
    [("The catch (be honest): ",MAGENTA,True),("measurement collapses the state and returns only n classical bits. You cannot read out all 2ⁿ amplitudes.",GRAY4,False)],
    [("The art of quantum algorithms is arranging interference so the useful answer is what you measure with high probability",INK,True)],
    [("This honesty buys credibility with a room full of scientists.",GRAY4,True)],
],size=12,gap=9)

# 3.3 encoding
s,ct=header(prs,"Section 3 · Quantum","Loading classical data into a quantum state")
enc=[("Amplitude encoding",BLUE,["Load a vector x into the amplitudes of the state","A d-dimensional vector needs only ⌈log₂ d⌉ qubits","This is the exponential compression that makes high-D inner products cheap in principle","Powerful for dimensionality reduction"]),
     ("Angle encoding",TEAL,["Embed features into the rotation angles of quantum gates","Continuous features influence the state directly","Ideal for shallow, NISQ-friendly parameterised circuits","The route the demo notebook uses"])]
cwd=IN(6.0);ch=IN(2.6)
for i,(h,c,items) in enumerate(enc):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=15,bsize=11.5)
tb,tf=txt(s,IN(0.6),ct+IN(2.9),IN(12.2),IN(1.3))
rich(tf,[("Other schemes exist (basis, block, one-hot), but the two above carry the whole story: ",GRAY4,False),
         ("compress the data, then compare states.",INK,True)],12.5,first=True,after=0,line=1.2)

# 3.4 feature maps
s,ct=header(prs,"Section 3 · Quantum","Quantum feature maps — reach a space hard to fake classically")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(6.1),IN(4.2),[
    [("A quantum feature map φ encodes each point x into a quantum state |φ(x)⟩ via a parameterised circuit",GRAY4,False)],
    [("It projects data into an exponentially high-dimensional Hilbert space",INK,True)],
    [("Patterns impossible to separate classically can become linearly separable there",GRAY4,False)],
    [("Highly-entangled SU(2) / ZZ feature maps improve cluster stability",GRAY4,False)],
    [("Benchmarks on shallow, near-term circuits:",INK,True)],
],size=12,gap=9)
stat(s,IN(7.0),ct+IN(0.5),IN(2.75),IN(1.6),"88.6%","quantum-kernel clustering accuracy on Iris (SU(2) feature map)",color=BLUE,nsize=30)
stat(s,IN(10.0),ct+IN(0.5),IN(2.75),IN(1.6),"91.0%","and on the Breast-Cancer dataset — on shallow circuits",color=TEAL,nsize=30)
tb,tf=txt(s,IN(7.0),ct+IN(2.45),IN(5.75),IN(1.6))
para(tf,"The bet",13,INK,bold=True,first=True,after=4)
para(tf,"The quantum feature map reaches a Hilbert space that is expensive to simulate classically — so structure invisible to classical kernels becomes separable. Shallow, NISQ-friendly, runs today.",11.5,GRAY4,after=0,line=1.2)

# 3.5 inner product = similarity
s,ct=header(prs,"Section 3 · Quantum","Inner product = similarity: the whole game")
tb,tf=txt(s,IN(0.6),ct+IN(0.2),IN(12.0),IN(4.0))
para(tf,"The natural quantity a quantum computer gives you between two states |a⟩ and |b⟩ is their overlap ⟨a|b⟩.",14,INK,bold=True,first=True,after=10,line=1.2)
para(tf,"‖a − b‖²  =  2 − 2·⟨a|b⟩      (for normalised vectors)",16,BLUE,bold=True,font=MONO,after=10)
rich(tf,[("Inner products ",GRAY4,False),("are",INK,True),(" similarities. Distance is one step away from overlap.",GRAY4,False)],13.5,after=10,line=1.2)
rich(tf,[("So we need one small circuit that turns an overlap into a measurable probability. ",GRAY4,False),
         ("That circuit is the swap test — next slide.",MAGENTA,True)],13.5,after=0,line=1.2)

# 3.6 SWAP TEST (hero)
s,ct=header(prs,"Section 3 · Quantum","The Swap Test — similarity in one measurement")
pic(s,g("swap_test_circuit.png"),IN(0.5),ct+IN(0.05),IN(7.6),IN(3.9),valign="top")
tf=bullets(s,IN(8.2),ct+IN(0.2),IN(4.6),IN(4.0),[
    [("Prepare |a⟩, |b⟩ and an ancilla in |0⟩",GRAY4,False)],
    [("Hadamard the ancilla → superposition",GRAY4,False)],
    [("Controlled-SWAP the two registers",GRAY4,False)],
    [("Hadamard the ancilla again; measure it",GRAY4,False)],
    [("P(ancilla=0) = ½ + ½·|⟨a|b⟩|²",BLUE,True,MONO)],
    [("Similar → P₀ near 1 · orthogonal → P₀ = ½",INK,True)],
    [("Measures the similarity of two d-dimensional states with a constant-size circuit — the seed of Lloyd–Mohseni–Rebentrost quantum k-means (2013)",MAGENTA,True)],
],size=11.5,gap=7)

# 3.7 q-means assemble + table
s,ct=header(prs,"Section 3 · Quantum","Quantum k-means (q-means) — assemble the pieces")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(5.3),IN(2.4),[
    [("Same Lloyd's loop — but the assign step (the n×k distances) uses the swap test on amplitude-encoded points",GRAY4,False)],
    [("Under ideal data loading (QRAM), per-distance cost scales like O(log(nd)) instead of O(nd)",INK,True)],
    [("Framework by Kerenidis, Landman, Luongo & Prakash — q-means++ init, quantum distance estimation",GRAY4,False)],
    [("Refined bounds give a quadratic runtime improvement over the best classical δ-k-means",GRAY4,False)],
],size=11,gap=6)
build_table(s,IN(6.1),ct+IN(0.15),IN(6.7),IN(3.4),
    [["","Classical k-means","Quantum q-means"],
     ["Distance","sequential Euclidean","swap-test inner product"],
     ["Feature space","original dimensionality","exponential Hilbert space"],
     ["Scaling in n","strictly linear O(n)","sub-linear / logarithmic*"],
     ["Non-linear","needs heavy Gram matrix","native via SU(2)/ZZ maps"],
     ["Hardware","CPU / GPU","QPU (+ QRAM for max speedup)"]],
    colw=[IN(1.5),IN(2.6),IN(2.6)],header_fill=PURPLE,fs=10)
tb,tf=txt(s,IN(6.1),ct+IN(3.65),IN(6.7),IN(0.5))
para(tf,"*logarithmic scaling assumes efficient QRAM — see the reality check.",9.5,GRAY,italic=True,first=True,after=0)

# 3.8 quantum kernel clustering
s,ct=header(prs,"Section 3 · Quantum","Quantum-kernel clustering — the NISQ-friendly route")
pic(s,g("kernel_heatmap.png"),IN(8.4),ct+IN(0.1),IN(4.4),IN(3.7),valign="top")
tb,tf=txt(s,IN(0.6),ct+IN(0.1),IN(7.5),IN(1.2))
para(tf,"Don't speed up k-means — change the similarity to one expensive to fake classically.",13.5,INK,bold=True,first=True,after=4,line=1.15)
para(tf,"K(xᵢ, xⱼ) = |⟨φ(xᵢ)|φ(xⱼ)⟩|²",15,BLUE,bold=True,font=MONO,after=0)
flow(s,IN(0.6),ct+IN(1.55),IN(7.5),IN(0.95),
     ["Data\nx","Quantum\nfeature map\n|φ(x)⟩","Kernel\nmatrix K\n(swap test)","Spectral\nclustering"],
     color=PURPLE,fs=10.5,agap=0.18)
tf=bullets(s,IN(0.6),ct+IN(2.75),IN(7.5),IN(1.6),[
    [("Feed K into any kernel method you already teach — spectral clustering, kernel k-means",GRAY4,False)],
    [("This is the near-term, runnable approach — shallow, NISQ-tolerant — and it closes the loop with Section 2's spectral clustering",TEAL,True)],
    [("The heatmap shows the block structure a good quantum kernel exposes",GRAY4,False)],
],size=11,gap=6)

# 3.9 DEMO
s,ct=header(prs,"Section 3 · Quantum","Live demo — k-means fails on the moons; the quantum kernel doesn't")
pic(s,g("moons_result.png"),IN(0.6),ct+IN(0.1),IN(8.1),IN(3.9),valign="top")
tf=bullets(s,IN(8.9),ct+IN(0.2),IN(3.9),IN(4.0),[
    [("qml_clustering_demo.ipynb",BLUE,True,MONO)],
    [("Two interleaved moons, 2-qubit data-reuploading feature map, kNN-sparsified, fed to SpectralClustering",GRAY4,False)],
    [("Classical k-means: ARI ≈ 0.48 — slices straight through the crescents",INK,True)],
    [("Quantum-kernel spectral: ARI ≈ 0.67 — recovers them",TEAL,True)],
    [("2 qubits, runs on this laptop — no quantum hardware",MAGENTA,True)],
    [("Teaching honesty: a curated example where the geometry favours the quantum kernel — say so, and always benchmark against a strong classical baseline",GRAY4,False)],
],size=11,gap=7)

# 3.10 own work
s,ct=header(prs,"Section 3 · Quantum","From my own work — the same circuits, put to work on portfolios")
pic(s,g("vqe_circuit.png"),IN(0.5),ct+IN(0.1),IN(7.2),IN(2.9),valign="top")
tf=bullets(s,IN(7.9),ct+IN(0.15),IN(4.9),IN(3.0),[
    [("Quantum Finance — An Overview (EasyChair 6071, 2021)",INK,True)],
    [("Mean-variance portfolio optimization → QUBO → Ising Hamiltonian → solved with a VQE",GRAY4,False)],
    [("Implemented in BOTH Qiskit and PennyLane — the framework choice this room faces",GRAY4,False)],
    [("Qiskit: optimal selection [0 1 0 1] @ prob 0.94; PennyLane VQE → [1 0 1]",INK,True)],
],size=11,gap=6)
tb,tf=txt(s,IN(0.6),ct+IN(3.15),IN(12.2),IN(1.3))
rich(tf,[("The bridge:  ",MAGENTA,True),
         ("the VQE ansatz (RX/RY + CNOT) is the same species of circuit as today's quantum feature map. There we optimised it to minimise a Hamiltonian; here we evaluate it to measure similarity. Same physics, two jobs.",GRAY4,False)],13,first=True,after=0,line=1.25)

# 3.11 reality check
s,ct=header(prs,"Section 3 · Quantum","Reality check — what is real, and what waits on hardware")
pic(s,g("advantage_quadrant.png"),IN(8.2),ct+IN(0.1),IN(4.6),IN(3.9),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(7.3),IN(4.2),[
    [("Data loading is the bottleneck. ",MAGENTA,True),("The exponential k-means speedup assumes QRAM to load classical data in superposition — practical QRAM does not yet exist at scale. State the assumption every time.",GRAY4,False)],
    [("Dequantization (Ewin Tang, 2018+). ",MAGENTA,True),("Several 'exponential' QML speedups were matched classically under the same sampling assumptions.",GRAY4,False)],
    [("NISQ noise → keep circuits shallow (hence shallow kernels). Barren plateaus → deep variational maps can have vanishing gradients.",GRAY4,False)],
    [("Where the win is genuine: ",TEAL,True),("quantum-native / structured data, small-but-hard kernels, and as a feature generator feeding the classical methods you already teach.",GRAY4,False)],
    [("Not 'k-means but faster on your CSV' — a new kind of similarity you can run now.",INK,True)],
],size=11.5,gap=8)

# =====================================================================
# SECTION 4 — INDUSTRY + MEGA-SCIENCE
# =====================================================================
divider(prs,"04","Industry & Current Research","From commercial pilots to the world's largest physics experiments — where quantum clustering is running now.",MAGENTA)

# 4.1 turning point
s,ct=header(prs,"Section 4 · Industry & Research","2025–26: from lab curiosity to commercial pilots")
tiles=[("300+","enterprises now engaged in quantum (McKinsey 2026)",BLUE),
       ("6.3×","jump in start-up investment, 2024 → 2025",TEAL),
       ("~90%","of that investment flowing to quantum computing",PURPLE),
       ("$4.4B","projected quantum-computing revenue by 2028",MAGENTA)]
tw=IN(2.92);th=IN(1.7)
for i,(n,l,c) in enumerate(tiles):
    stat(s,IN(0.57)+i*(tw+IN(0.16)),ct+IN(0.1),tw,th,n,l,color=c,nsize=27)
tb,tf=txt(s,IN(0.6),ct+IN(2.1),IN(12.2),IN(2.4))
para(tf,"The discourse shifted from 'if' quantum advantage would occur to modelling exactly 'when' it disrupts each sector.",13.5,INK,bold=True,first=True,after=8,line=1.2)
para(tf,"The pattern across every application below is the same: high-dimensional, mostly unlabelled data, where the expensive step is computing similarities — the corner quantum targets. Finance, pharma and climate are the commercial vanguard; the mega-science projects are where the hardest versions run today.",12.5,GRAY4,after=0,line=1.25)

# 4.2 finance
s,ct=header(prs,"Section 4 · Industry & Research","Financial services — the early adopters")
fin=[("Option pricing & risk",BLUE,["JPMorgan Chase × IBM Quantum","Iterative Quantum Amplitude Estimation (IQAE)","Quadratic speedup over classical Monte Carlo — error ε needs 1/ε samples, not 1/ε²","Fixed confidence far faster, or higher confidence at fixed time"]),
     ("Fraud detection",TEAL,["Quantum SVM / QSVM with ZZ-feature-maps (IBM Safer Payments, Qiskit)","Projects transactions where fraud & legitimate behaviour separate cleanly","VQC F1-scores up to 0.88, robust across 5 simulated noise types","Hybrid: QSVM for feature selection → classical routing"]),
     ("Deep hedging",PURPLE,["JPMorgan × QC Ware","Dynamic risk mitigation that adapts to market volatility in real time","Portfolio optimization = unsupervised structure-finding on a covariance matrix"])]
cwd=IN(3.95);ch=IN(3.1)
for i,(h,c,items) in enumerate(fin):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=13.5,bsize=10.5)

# 4.3 pharma + genomics
s,ct=header(prs,"Section 4 · Industry & Research","Pharmaceuticals, healthcare & genomics")
stat(s,IN(0.57),ct+IN(0.1),IN(3.0),IN(1.7),"97.3%","Dynamic Quantum Clustering of RNA-seq (20,057 genes) → 90-gene panel, no labels",color=TEAL,nsize=28)
stat(s,IN(3.75),ct+IN(0.1),IN(3.0),IN(1.7),"8-qubit","Algorithmiq × Quantum Circuits 'Aqumen Seeker' dual-rail QPU for enzyme pharmacokinetics",color=BLUE,nsize=24)
stat(s,IN(6.93),ct+IN(0.1),IN(3.0),IN(1.7),"90.9%","DQC concordance with clinical glioma outcomes — geometry-preserving, unsupervised",color=PURPLE,nsize=28)
stat(s,IN(10.1),ct+IN(0.1),IN(2.7),IN(1.7),"de novo","structure-based drug design — 'correct first, then scale' error detection",color=MAGENTA,nsize=24)
tf=bullets(s,IN(0.6),ct+IN(2.05),IN(12.2),IN(2.3),[
    [("Genomics is the ultimate clustering challenge: ",INK,True),("phenotypic significance among tens of thousands of correlated genes, with no preset labels or cluster counts.",GRAY4,False)],
    [("Dynamic Quantum Clustering (DQC) organised bulk RNA-seq into clusters that mirrored the LGG→GBM disease continuum, then a 90-gene panel lifted accuracy to 97.3%",GRAY4,False)],
    [("The biochemical realm is natively quantum — classical simulation of molecules fails exponentially with size; dual-rail qubits detect errors as photon-loss 'erasures' and post-select",GRAY4,False)],
],size=11.5,gap=8)

# 4.4 climate + energy
s,ct=header(prs,"Section 4 · Industry & Research","Sustainability, climate & energy")
cl=[("Carbon capture (MOFs)",TEAL,["Quantinuum × TotalEnergies","Quantum fragmentation models CO₂ binding in Metal-Organic Frameworks — 'molecular LEGO'","Precise many-body binding energies, avoiding trial-and-error delays","PrISMa tool blends ML + quantum-informed simulation"]),
    ("Catalysis",BLUE,["Simulate N₂ dissociation on a platinum surface with quantum precision","Ammonia-synthesis catalysts ≈ 2% of global CO₂ — push efficiency toward 100%","Pulse-electrolysis energy states, catalyst-membrane interactions"]),
    ("Grid & sensing",PURPLE,["Quantum k-means (angle + amplitude embeddings) on the German electricity grid","+67.8% balanced accuracy for predictive maintenance & grid labelling","Cold-atom gravity sensors monitor underground CO₂ storage plumes"])]
cwd=IN(3.95);ch=IN(3.1)
for i,(h,c,items) in enumerate(cl):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=13.5,bsize=10.3)

# 4.5 mega-science intro
s=prs.slides.add_slide(L(prs,"Blank slide"))
rect(s,0,0,SW,SH,INK)
tb,tf=txt(s,IN(0.9),IN(0.9),IN(11.5),IN(1.6))
para(tf,"MEGA-SCIENCE",13,CY,bold=True,first=True,after=6)
para(tf,"The biggest experiments ever built are unsupervised-learning machines.",30,WHITE,bold=True,after=0,line=1.05)
pic(s,g("data_deluge.png"),IN(8.0),IN(2.7),IN(4.9),IN(3.4),valign="top")
tf=bullets(s,IN(0.9),IN(2.8),IN(6.7),IN(3.6),[
    [("SKA streams ~8 Tbit/s and will archive ~8.5 exabytes",WHITE,True)],
    [("The HL-LHC (2029–42) collects >10× today's data, into the exabyte regime, needing ~50–100× the compute",LGRAY,False)],
    [("DUNE turns each neutrino into a 3D argon image — 30–60 PB/year",LGRAY,False)],
    [("LIGO reads a 4-km laser ruler perturbed by 10⁻¹⁹ m",LGRAY,False)],
    [("All of it: extreme-dimensional, mostly UNLABELLED data — clustering / anomaly-detection / unsupervised-representation territory",CY,True)],
],size=12,gap=9)
pageno(s)

# 4.6 LIGO
s,ct=header(prs,"Section 4 · Mega-Science","LIGO — quantum-limited sensing meets ML")
pic(s,g("ligo_squeeze.png"),IN(7.3),ct+IN(0.1),IN(5.5),IN(3.5),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(6.5),IN(4.2),[
    [("Frequency-dependent squeezing (A+): ",INK,True),("4.0 dB (Hanford) / 5.8 dB (Livingston) of quantum-noise reduction via new 300 m filter cavities",GRAY4,False)],
    [("Extends detector range 15–18% → up to +65% detection rate in the O4 run",TEAL,True)],
    [("Gravity Spy (citizen science + CNN): ",INK,True),("23 glitch classes, >30k volunteers, >7M classifications; new multi-view CNN hits 94.1% accuracy / AUC 0.965 — unsupervised-flavoured glitch clustering",GRAY4,False)],
    [("Quantum ML on GW data: ",INK,True),("'quantum variational rewinding' treats readings as anomalies in background noise — linear-time vs quasilinear matched filtering",GRAY4,False)],
    [("LIGO-India broke ground 23 Apr 2026 (Aundha, Maharashtra) — two 4 km arms, ~2030, ≈₹2,600 cr",MAGENTA,True)],
],size=11,gap=7)

# 4.7 CERN
s,ct=header(prs,"Section 4 · Mega-Science","CERN / LHC — quantum clustering of collision events")
pic(s,g("cern_clustering.png"),IN(8.5),ct+IN(0.1),IN(4.3),IN(3.9),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(7.6),IN(4.2),[
    [("CERN Quantum Technology Initiative (est. 2020) — four Centres of Competence spanning QC, sensing & networks",GRAY4,False)],
    [("The headline unsupervised win: ",MAGENTA,True),("quantum anomaly detection in the LHC latent space — QK-means / QK-medians + a quantum kernel on autoencoder features found a regime where the quantum model significantly outperforms classical on real hardware (Commun. Phys. 7, 334, 2024)",GRAY4,False)],
    [("Track reconstruction cast as a QUBO on a D-Wave annealer (33 qubits, ~500 tracks); annealing-inspired classical then gave ~4-orders speed-up toward HL-LHC densities",GRAY4,False)],
    [("QSVM for Higgs (ttH): matched classical SVM/BDT at up to 20 qubits / 50,000 events",GRAY4,False)],
    [("Driver: HL-LHC storage enters the exabyte regime; compute ~50–100× today's",INK,True)],
],size=11,gap=6)

# 4.8 ITER
s,ct=header(prs,"Section 4 · Mega-Science","ITER & fusion — clustering plasma states, dodging disruptions")
pic(s,g("iter_plasma.png"),IN(8.2),ct+IN(0.1),IN(4.6),IN(3.6),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(7.3),IN(4.2),[
    [("DeepMind × EPFL (Nature 2022): ",INK,True),("deep-RL directly controlled a tokamak's magnetic coils — first RL magnetic control on a real machine",GRAY4,False)],
    [("DIII-D (Nature 2024): ",INK,True),("an RL controller forecasts tearing instabilities up to 300 ms ahead and steers actuators to stay stable",GRAY4,False)],
    [("Cross-machine disruption prediction (FRNN, Nature 2019): trained on DIII-D+JET, predicts on an untrained machine — crucial for ITER's ~30 ms warning budget",GRAY4,False)],
    [("Unsupervised classification separates confinement regimes (L-mode, H-mode, disruptive) in diagnostic space",TEAL,True)],
    [("Quantum: algorithms for plasma turbulence/transport; NV-center diamond magnetometry (to ~1.2 T) for radiation-hard diagnostics. ITER new baseline: Q ≥ 10 by 2044",MAGENTA,True)],
],size=11,gap=6)

# 4.9 SKA + DUNE
s,ct=header(prs,"Section 4 · Mega-Science","SKA & DUNE — label-free pipelines on petascale data")
card(s,IN(0.57),ct+IN(0.1),IN(6.0),IN(4.1),"SKA — radio astronomy",[
    "~8 Tbit/s to the correlators; ~135 PFLOPS each; >700 PB/yr → ~8.5 EB archive; >100× current internet traffic",
    "Convolutional-autoencoder + nearest-latent-neighbour flags RFI without labels; DBSCAN groups false positives",
    "Anomaly detection + active learning on MeerKAT surfaces rare transients",
    "Quantum-kernel SVM for galaxy morphology: ROC AUC 0.946 — at parity with classical",
    "First image Mar 2025: 85 galaxies in 25 deg² with <1% of antennas"],accent=BLUE,tsize=14,bsize=10.3)
card(s,IN(6.8),ct+IN(0.1),IN(6.0),IN(4.1),"DUNE — neutrinos",[
    "4 liquid-argon TPC modules, 70 kt, ~1.5 km underground; 30–60 PB/year",
    "Each event is a high-res 3D image → CNN / graph-net reconstruction; clustering separates overlapping tracks & showers",
    "Neural Projected Quantum Kernels + QCNN classify track vs cascade at ~80% (1–100 TeV), validated on IBM 127-qubit hardware",
    "VQC on LHC data reached AUC 0.81 — near-parity with classical",
    "Cryostat steel lowered underground in 2026; first modules 2026–27"],accent=TEAL,tsize=14,bsize=10.3)

# 4.10 quantum sensing + synthesis
s,ct=header(prs,"Section 4 · Mega-Science","Quantum sensing — a new class of instruments (and data)")
pic(s,g("sensing_sensitivity.png"),IN(7.2),ct+IN(0.1),IN(5.6),IN(3.2),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(6.4),IN(3.4),[
    [("Cold-atom gravimeter: 2.2 µGal·Hz⁻¹ᐟ²; Birmingham's gradiometer made the first outdoor detection of a buried tunnel — ~10× faster gravity surveys",GRAY4,False)],
    [("Optical clocks: 8.1×10⁻¹⁹ systematic uncertainty (JILA) — cm-level chronometric levelling, tests of relativity",GRAY4,False)],
    [("SERF magnetometers: ~4.5 fT·Hz⁻¹ᐟ², surpassing SQUIDs",GRAY4,False)],
    [("Km-scale atom interferometers — MAGIS-100, AION, AEDGE — hunt gravitational waves & dark matter in the mid-band between LISA and LIGO",GRAY4,False)],
],size=11,gap=7)
tb,tf=txt(s,IN(0.6),ct+IN(3.55),IN(12.2),IN(1.4))
rich(tf,[("The synthesis:  ",MAGENTA,True),
         ("mega-science emits extreme-dimensional, unlabelled data — clustering's home turf. Quantum helps from both ends: kernels/annealers embed the data (mostly parity today, but a real win for CERN's unsupervised latent-space clustering), and quantum sensors generate a new class of ultra-precise measurements that themselves become unlabelled streams.",GRAY4,False)],12,first=True,after=0,line=1.22)

# =====================================================================
# SECTION 5 — INDIA
# =====================================================================
divider(prs,"05","The Indian Landscape","India is building the hardware. This room builds the people who will use it.",DEEP)

# 5.1 NQM
s,ct=header(prs,"Section 5 · India","The National Quantum Mission")
pic(s,g("nqm_targets.png"),IN(8.4),ct+IN(0.1),IN(4.4),IN(3.3),valign="top")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(7.5),IN(2.6),[
    [("₹6,003.65 crore, approved 19 Apr 2023, running 2023–2031 (~$0.73B)",INK,True)],
    [("Staged compute: 20–50 qubits (yr 3) → 50–100 (yr 5) → 50–1000 (yr 8), superconducting + photonic",GRAY4,False)],
    [("152 researchers · 43 institutions · 17 states — a genuinely national footprint relevant to every college in the room",GRAY4,False)],
],size=11.5,gap=7)
hubs=[("Computing","IISc Bengaluru",BLUE),("Communication","IIT Madras + C-DOT",TEAL),
      ("Sensing & Metrology","IIT Bombay",PURPLE),("Materials & Devices","IIT Delhi",MAGENTA)]
hw=IN(2.95);hh=IN(1.1)
for i,(t,inst,c) in enumerate(hubs):
    x=IN(0.57)+i*(hw+IN(0.15))
    b=rect(s,x,ct+IN(2.85),hw,hh,PANEL,line=LGRAY,lw=0.75,round=True,rad=0.06);rect(s,x,ct+IN(2.85),hw,IN(0.1),c)
    tbb,tff=txt(s,x+IN(0.18),ct+IN(3.05),hw-IN(0.3),hh-IN(0.3))
    para(tff,t,12.5,INK,bold=True,first=True,after=2);para(tff,inst,10.5,GRAY4,after=0)
tb,tf=txt(s,IN(0.6),ct+IN(4.05),IN(12),IN(0.4))
para(tf,"Four thematic hubs, hub-and-spoke, each a Section-8 company — operational since Oct 2024.",10.5,GRAY,italic=True,first=True,after=0)

# 5.2 ecosystem
s,ct=header(prs,"Section 5 · India","Indigenous ecosystem — shipping, not slideware")
eco=[("QpiAI (Bengaluru)",BLUE,["'Indus' 25-qubit full-stack (World Quantum Day, Apr 2025)","'Kaveri' 64-qubit chip (Nov 2025), commercial late 2026","$65.6M raised, roadmap to 1,000 qubits by 2030"]),
     ("Amaravati Quantum Valley",TEAL,["Launched 7 Feb 2026 (IBM + TCS + Andhra Pradesh)","IBM Quantum System Two — 156-qubit Heron, India's largest","Commissioning ~Sept 2026; L&T infrastructure"]),
     ("Academic + startups",PURPLE,["IISER Pune ion-trap: 20-qubit ⁴⁰Ca⁺ machine due end-2026","~1,000 km quantum comms (QNu Labs ARMOS QKD, Apr 2026)","8→17 NQM-funded startups: QNu, Dimira, Quanastra, Prenishq…"])]
cwd=IN(3.95);ch=IN(2.7)
for i,(h,c,items) in enumerate(eco):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=13,bsize=10.3)
tb,tf=txt(s,IN(0.6),ct+IN(3.0),IN(12.2),IN(1.2))
rich(tf,[("Also 2026: ",INK,True),("a national post-quantum-cryptography roadmap (Feb 2026) to protect Aadhaar-scale identity and UPI-scale payments before the threat matures. Minister Jitendra Singh: India has met roughly half of NQM's targets within three years.",GRAY4,False)],12,first=True,after=0,line=1.2)

# 5.3 funding honesty
s,ct=header(prs,"Section 5 · India","The honest number — and where India's edge really is")
pic(s,g("funding_bars.png"),IN(0.5),ct+IN(0.1),IN(7.0),IN(3.9),valign="top")
tf=bullets(s,IN(7.7),ct+IN(0.3),IN(5.1),IN(4.0),[
    [("India's ~$0.73B over 8 years sits against China ~$15B, US $2.5B+, UK £2.5B/10yr, Germany ~€2–3B",GRAY4,False)],
    [("India is funded to compete selectively in chosen niches — not to outspend anyone",INK,True)],
    [("Its edge is algorithms, software, and a vast talent pool — the value-capturing layer",TEAL,True)],
    [("The winning move is the algorithm-and-application layer — which is exactly what QML, clustering and quantum-enhanced vision are",MAGENTA,True)],
    [("Faculty respect candor — say the number out loud.",GRAY4,True)],
],size=11.5,gap=8)

# 5.4 Viksit Bharat
s,ct=header(prs,"Section 5 · India","Quantum is a pillar of Viksit Bharat 2047")
tf=bullets(s,IN(0.6),ct+IN(0.1),IN(7.2),IN(4.2),[
    [("Viksit Bharat @2047: ",INK,True),("a fully developed India by the centenary of independence — prosperity, sustainability, effective governance",GRAY4,False)],
    [("India already took Aadhaar, UPI and Digital India to population scale — proven ability to run frontier tech at planet scale",GRAY4,False)],
    [("The next leap: 'world's back office' → 'world's innovation engine' — deep R&D and sovereign tech ecosystems",INK,True)],
    [("A coordinated mission stack: NQM + IndiaAI Mission + National Supercomputing Mission + Semicon India + ANRF (₹1-lakh-crore RDI fund)",GRAY4,False)],
    [("Sovereignty: secure comms, indigenous hardware — not renting foreign quantum cloud",TEAL,True)],
],size=11.5,gap=8)
card(s,IN(8.1),ct+IN(0.1),IN(4.7),IN(4.1),"The gap to close (be honest)",[
    "R&D spend must rise toward 2–3% of GDP by 2047",
    "Reverse the brain drain; create high-tech jobs",
    "Every faculty member here is a lever on the talent side of that equation",
    "\"Viksit Bharat 2047 will be built by the cohort sitting in your classrooms right now.\""],accent=DEEP,tsize=13,bsize=11)

# =====================================================================
# SECTION 6 — FUTURE
# =====================================================================
divider(prs,"06","The Future","The single largest bottleneck globally is people who understand this stack. You are the lever.",INK)

# 6.1 where the field goes
s,ct=header(prs,"Section 6 · The Future","Where the field goes — fault-tolerance & AI × Quantum")
fut=[("Fault-tolerant roadmaps",BLUE,["IBM Starling (2029): 200 logical qubits, 100M gates; Blue Jay → 2,000 logical","Google: useful error-corrected machine ~2029, end-state ~1M physical qubits","Quantinuum Helios → Apollo (2029); PsiQuantum photonic million-qubit"]),
     ("AI × Quantum",TEAL,["AlphaQubit (DeepMind, Nature 2024): neural decoder for the surface code, 6–30% better error decoding","AI decoders rescue noisy NISQ hardware — the two fields have merged","Quantum kernels/feature maps for structured data"]),
     ("The data-loading frontier",PURPLE,["QRAM is the Achilles heel of QML speedups","Zhejiang: first superconducting bucket-brigade QRAM demo (Jun 2026), 4/8-bit addressing","Progress here unlocks the exponential clustering speedups"])]
cwd=IN(3.95);ch=IN(3.2)
for i,(h,c,items) in enumerate(fut):
    card(s,IN(0.57)+i*(cwd+IN(0.2)),ct+IN(0.1),cwd,ch,h,items,accent=c,tsize=13.5,bsize=10.5)

# 6.2 talent + call to action
s,ct=header(prs,"Section 6 · The Future","The bottleneck is people — you are the lever")
pic(s,g("talent_gap.png"),IN(0.5),ct+IN(0.1),IN(6.7),IN(3.9),valign="top")
card(s,IN(7.4),ct+IN(0.1),IN(5.4),IN(3.9),"Call to action for faculty",[
    "Clustering, kernels and the swap test are teachable from first principles — no hardware required",
    "Add one quantum-ML module to an existing ML course",
    "Use these notebooks — they run on any laptop in PennyLane's CPU simulator",
    "Point strong students at the NQM hubs and QpiAI / BosonQ / startup internships",
    "Every course you seed, every student you point at QML, is India's comparative advantage compounding"],
    accent=BLUE,tsize=14,bsize=11.5)

# 6.3 THANK YOU
s=prs.slides.add_slide(L(prs,"Blank slide"))
rect(s,0,0,SW,SH,INK)
rect(s,0,0,IN(0.28),SH,BLUE)
tb,tf=txt(s,IN(1.0),IN(1.1),IN(11),IN(2.2))
para(tf,"Thank you.",46,WHITE,bold=True,first=True,after=8)
para(tf,"Clustering, reimagined — from first principles, no black boxes. Bring your students to the next cohort.",15,LGRAY,after=0,line=1.2)
# contact block
tb,tf=txt(s,IN(1.0),IN(3.6),IN(7.4),IN(3.0))
rich(tf,[("Archit Srivastava",WHITE,True)],20,first=True,after=6)
rich(tf,[("Founder, AiQyaM — India's first quantum-hardware community",LGRAY,False)],12.5,after=10)
rich(tf,[("Email   ",CY,True),("architsrivastava3115@gmail.com",WHITE,False)],13,after=4)
rich(tf,[("Scholar  ",CY,True),("NbPUdWMAAAAJ  (scholar.google.com)",WHITE,False)],13,after=4)
rich(tf,[("Session ",CY,True),("Quantum Unsupervised Learning (Clustering) · QML-2026 · 11 Jul 2026",WHITE,False)],13,after=4)
rich(tf,[("Notebooks  ",CY,True),("qml_clustering_demo.ipynb — runs on any laptop, no quantum hardware",WHITE,False)],13,after=0)
try: s.shapes.add_picture(a("logo-wordmark-light.png"),IN(9.5),IN(6.35),height=IN(0.6))
except Exception: pass
PAGENO[0]+=1

prs.save(OUT)
print("SAVED",OUT)
print("SLIDES",len(prs.slides._sldIdLst))
